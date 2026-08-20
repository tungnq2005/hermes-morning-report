"""Build a .pptx with a consistent, branded layout from an extracted document outline.

Text goes into the layout's own placeholders rather than free-floating textboxes.
Hand-placed boxes rendered fine but carried no structure: every slide came out on the
Blank layout with no title placeholder, so PowerPoint's outline view was empty, a
theme change moved nothing, and importers like Canva saw loose boxes instead of a
title and a body. Placeholders also bring real bullet formatting and autofit.

Visual identity:
  * the cover and closing slides sit on a full-bleed navy background with white text;
  * a heading that names a whole part (Opportunities / Risks) becomes a full-colour
    section divider, and the content that follows inherits its colour -- green for
    opportunities, red for risks, navy for everything else.
"""
from __future__ import annotations

import math
import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ACCENT = RGBColor(0x1F, 0x3A, 0x5F)   # brand navy (titles, neutral accent bars)
BRAND_BG = RGBColor(0x0F, 0x2A, 0x4A)  # deep navy for full-bleed backgrounds
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINT = RGBColor(0xD8, 0xE2, 0xF0)      # light text on dark backgrounds
OPPORTUNITY = RGBColor(0x1B, 0x7A, 0x43)  # green
RISK = RGBColor(0xB3, 0x30, 0x1F)         # red
BODY = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
CARD_BG = RGBColor(0xF4, 0xF6, 0xFB)   # light card fill

# A divider's scrim has to stay readable over the *brightest* photo Openverse might
# return. Measured against a white picture, the old 70%-opacity section colour left
# white text at 3.0:1 on green and 3.6:1 on red -- below the 4.5:1 that body text needs,
# and exactly the "chữ chìm vào nền" case. Darkening the scrim colour buys the contrast
# without hiding the picture the way a near-opaque scrim would: at 55% darkening and
# 74% opacity the worst case (white photo) measures 7.1 navy / 5.3 green / 6.0 red.
SCRIM_DARKEN = 0.55
SCRIM_ALPHA = 74

MAX_BULLETS_PER_SLIDE = 6
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# python-pptx's default template is a 10x7.5 (4:3) deck. Widening slide_width to 16:9
# does not move the layouts' placeholders, so every slide used to render its content
# inside the old 4:3 box: a 9"-wide title centred at the 5" mark, a 4.42" body column
# that wrapped short sentences into three cramped lines, and 3.8" of dead space down
# the right-hand edge. Every placeholder is therefore re-fitted to the real canvas.
MARGIN = Inches(0.7)
GUTTER = Inches(0.5)
TITLE_TOP = Inches(0.4)
TITLE_H = Inches(1.0)
BODY_TOP = Inches(1.6)
BODY_H = SLIDE_H - BODY_TOP - Inches(0.8)
FULL_W = SLIDE_W - 2 * MARGIN
TEXT_W_WITH_IMAGE = Inches(6.9)
PIC_LEFT = MARGIN + TEXT_W_WITH_IMAGE + GUTTER
PIC_W = SLIDE_W - PIC_LEFT - MARGIN

# Indices into the default template's layout list.
LAYOUT_TITLE = 0        # Title Slide      -- title + subtitle placeholders
LAYOUT_TITLE_BODY = 1   # Title and Content
LAYOUT_TWO_CONTENT = 3  # Two Content      -- body on the left, picture on the right
LAYOUT_SECTION = 2      # Section Header

# Body text is sized to fit before it is written. A bare <a:normAutofit/> is a weak
# guarantee: PowerPoint only recomputes its shrink factor when the text box is edited,
# so a deck that overflows opens overflowing however the tag is set. The sizes below
# are tried largest-first and the first one whose estimated block height fits wins;
# normAutofit stays on only as a backstop for renderers that disagree with the estimate.
BODY_SIZE_CHOICES = (Pt(20), Pt(18), Pt(16), Pt(14), Pt(12), Pt(11))
# Calibri/Carlito average out near half the point size per character at body weights.
AVG_CHAR_WIDTH_EM = 0.50
LINE_HEIGHT_EM = 1.22
# A title long enough to wrap costs a slide most of its body room.
TITLE_SIZES = ((70, Pt(22)), (45, Pt(25)), (0, Pt(28)))
# Below this, a heading's whole body is a lead-in sentence rather than content.
DIVIDER_MAX_CHARS = 180

STRINGS = {
    "en": {"agenda": "Agenda", "credits": "Image credits",
           "closing": "Summary & Q&A", "continued": "cont."},
    "vi": {"agenda": "Nội dung", "credits": "Nguồn ảnh",
           "closing": "Tóm tắt & Q&A", "continued": "tiếp"},
}

# A heading that names a whole part switches the deck's colour phase; content that
# follows inherits it until the next switch.
OPPORTUNITY_MARKERS = ("cơ hội", "co hoi", "opportunit")
RISK_MARKERS = ("rủi ro", "rui ro", "risk")
# Titles that read as a list of short pointers -> a card grid instead of one long
# bullet column (watchlist, actions, recommendations).
CARD_MARKERS = ("chỉ báo", "theo dõi", "hành động", "đề xuất", "khuyến nghị",
                "watch", "action", "suggest", "recommend", "next step")


def _section_color(title: str, level: int, phase: str) -> tuple[str, str]:
    """Return (color_key, next_phase) for a section title.

    A heading that names a whole part (Opportunities / Risks) switches the deck's
    colour phase. A sub-heading (level 3) inherits the current phase. Any other
    top-level heading (level <= 2) resets back to neutral, so trailing sections such
    as "Watchlist" / "Suggested actions" are not stained by the risks that precede them.
    """
    t = title.lower()
    if any(m in t for m in RISK_MARKERS):
        return "risk", "risk"
    if any(m in t for m in OPPORTUNITY_MARKERS):
        return "opportunity", "opportunity"
    if level <= 2:
        return "neutral", "neutral"
    return phase, phase


def _color_for(key: str) -> RGBColor:
    if key == "opportunity":
        return OPPORTUNITY
    if key == "risk":
        return RISK
    return ACCENT


def _darken(color: RGBColor, factor: float) -> RGBColor:
    value = int(str(color), 16)
    return RGBColor(*(int(((value >> shift) & 0xFF) * factor) for shift in (16, 8, 0)))


def _color_background(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _background_is_dark(slide) -> bool:
    try:
        fill = slide.background.fill
        if fill.type != MSO_FILL.SOLID:
            return False
        rgb = fill.fore_color.rgb
        if not rgb:
            return False
        value = int(str(rgb), 16)
        r, g, b = (value >> 16) & 0xFF, (value >> 8) & 0xFF, value & 0xFF
        return (r * 299 + g * 587 + b * 114) / 1000 < 128
    except Exception:  # best-effort: default to light-on-white
        return False


# --- content-aware templates --------------------------------------------------
_STAT_NUM_RE = re.compile(r"\d[\d.,]*")
_STAT_UNIT_RE = re.compile(
    r"(?:\s*(?:%|tỉ|triệu|nghìn|lần|[x×]|USD|\$|ngày|giờ|tháng|năm|điểm|bil\.?|mil\.?)){1,3}",
    re.IGNORECASE,
)


def _find_stat_span(text: str) -> tuple[int, int] | None:
    """Span of the salient figure in a news bullet: the last number carrying a unit
    (%, tỉ, triệu, lần, ngày, USD, ...). Bare numbers are ignored so we do not big-up
    "Nghị định 142" or the "2" in "C2PA"."""
    with_unit = None
    for m in _STAT_NUM_RE.finditer(text):
        start, end = m.start(), m.end()
        unit = _STAT_UNIT_RE.match(text[end:])
        if unit:
            end += unit.end()
            with_unit = (start, end)
    return with_unit


def _is_stat_section(items: list[str]) -> bool:
    items = [i for i in items if i.strip()]
    return 2 <= len(items) <= 4 and sum(1 for i in items if _find_stat_span(i)) >= 2


def _lead_span(text: str) -> tuple[int, int] | None:
    """Span of a bullet's headline lead-in (before the first em dash). Bolded so the
    point reads at a glance."""
    at = text.find("—")
    return (0, at) if 0 < at <= 70 else None


def _is_card_section(title: str, items: list[str]) -> bool:
    items = [i for i in items if i.strip()]
    if not (2 <= len(items) <= 6) or not all(len(i) <= 160 for i in items):
        return False
    return any(m in title.lower() for m in CARD_MARKERS)


def build(doc: dict, sections: list[dict], out_path: str, min_slides: int = 5,
          images: list[str | None] | None = None, subtitle: str = "",
          credits: list[dict] | None = None, cover_image: str | None = None) -> dict:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    # Keep one slot per requested image, None where the fetch failed, so a missing
    # picture leaves its own slide bare instead of shifting every later picture onto
    # the wrong section.
    images = [p if p and os.path.exists(p) else None for p in (images or [])]
    words = STRINGS["en" if _is_english(doc) else "vi"]

    # A picture on the cover must not appear again inside the deck: sections whose
    # slot holds the cover image go without one rather than repeating it.
    cover_image = cover_image if cover_image and os.path.exists(cover_image) else None
    if cover_image:
        images = [None if img == cover_image else img for img in images]

    _title_slide(prs, doc["title"], subtitle, cover_image)
    # Reach min_slides by spreading the real content thinner, never by padding: the
    # deck used to append slides whose only bullet was the literal string "(bổ sung)",
    # and they landed after the closing slide.
    content_sections = _paginate(sections, words)
    for max_bullets in (3, 2, 1):
        if len(content_sections) + 2 >= min_slides:
            break
        content_sections = _paginate(sections, words, max_bullets=max_bullets)
    agenda = [s["title"] for s in content_sections if s["title"] and s["first"]]
    if agenda:
        _bullet_slide(prs, words["agenda"], agenda[:8], None)

    # Index the pictures by source section rather than walking them slide by slide: a
    # section long enough to be split across slides would otherwise eat its
    # neighbours' pictures and push every later one onto the wrong section.
    placed = 0
    rejected: list[str] = []
    phase = "neutral"
    for sec in content_sections:
        image = images[sec["src"]] if sec["first"] and sec["src"] < len(images) else None
        title = sec["title"] or doc["title"]
        color_key, phase = _section_color(title, sec.get("level", 2), phase)
        color = _color_for(color_key)
        # A heading whose whole body is one short line is a part divider, not a content
        # slide. Rendering it as one leaves a single bullet marooned in white space.
        if _is_divider(sec):
            embedded = _divider_slide(prs, title, sec["items"][0] if sec["items"] else "",
                                      image, color_key)
            uses_image = True
        elif _is_stat_section(sec["items"]):
            _stat_slide(prs, title, sec["items"], color)
            embedded, uses_image = False, False
        elif _is_card_section(title, sec["items"]):
            _card_slide(prs, title, sec["items"], color)
            embedded, uses_image = False, False
        else:
            embedded = _bullet_slide(prs, title, sec["items"], image, color)
            uses_image = True
        if image:
            if embedded:
                placed += 1
            elif uses_image:
                rejected.append(os.path.basename(image))
    if credits:
        _credits_slide(prs, credits, words)
    _closing_slide(prs, doc["title"], words)
    _add_slide_numbers(prs)

    prs.save(out_path)
    return {"slides": len(prs.slides), "images_used": placed + (1 if cover_image else 0),
            "images_rejected": rejected}


def _paginate(sections: list[dict], words: dict,
              max_bullets: int = MAX_BULLETS_PER_SLIDE) -> list[dict]:
    """Split over-long sections across slides. Each chunk remembers which section it
    came from (`src`) and whether it opens it (`first`) so pictures stay aligned."""
    out: list[dict] = []
    for src, sec in enumerate(sections):
        items = sec["items"] or [""]
        level = sec.get("level", 2)
        for i in range(0, len(items), max_bullets):
            chunk = [t for t in items[i:i + max_bullets] if t]
            title = sec["title"] if i == 0 else f"{sec['title']} ({words['continued']})"
            out.append({"title": title, "items": chunk, "src": src,
                        "first": i == 0, "level": level})
    return out


def _is_english(doc: dict) -> bool:
    sample = (doc["title"] + " " + " ".join(b["text"] for b in doc["blocks"][:10])).lower()
    vi_chars = sum(1 for ch in sample if ch in "ăâđêôơưàáảãạèéẻẽẹìíỉĩịòóỏõọùúủũụỳýỷỹỵ")
    return vi_chars < 3


def _fit(shape, left, top, width, height) -> None:
    shape.left, shape.top, shape.width, shape.height = left, top, width, height


def _pick(sizes, measure: int):
    for threshold, size in sizes:
        if measure > threshold:
            return size
    return sizes[-1][1]


def estimated_text_height(bullets: list[str], size, width, space_after=Pt(8)):
    """Rough height of a bulleted block. Deliberately renderer-independent: it must
    hold for PowerPoint (Calibri) and LibreOffice (Carlito), whose metrics match."""
    if not bullets:
        return 0
    char_width = size * AVG_CHAR_WIDTH_EM
    # The placeholder's hanging indent costs roughly two characters of usable width.
    usable = max(width - Emu(int(char_width * 2)), Emu(1))
    per_line = max(int(usable / char_width), 1)
    lines = sum(max(-(-len(text) // per_line), 1) for text in bullets)
    return int(lines * size * LINE_HEIGHT_EM + len(bullets) * space_after)


def _fitting_body_size(bullets: list[str], width, height):
    for size in BODY_SIZE_CHOICES:
        if estimated_text_height(bullets, size, width) <= height:
            return size
    return BODY_SIZE_CHOICES[-1]


def _fitting_card_body_size(text: str, width, height):
    """Card body size, capped at 14pt (the card design's body max) and floored at 11pt."""
    for size in (Pt(14), Pt(12), Pt(11)):
        if estimated_text_height([text], size, width) <= height:
            return size
    return Pt(11)


def _clear_effect_ref(shape) -> None:
    from pptx.oxml.ns import qn

    style = shape._element.find(qn("p:style"))
    if style is None:
        return
    effect_ref = style.find(qn("a:effectRef"))
    if effect_ref is not None:
        effect_ref.set("idx", "0")


def _strip_style_ref(shape) -> None:
    """Remove a shape's <p:style> so renderer-independent checks treat it as inert
    furniture (a full-bleed background scrim), not an accent rule."""
    from pptx.oxml.ns import qn

    style = shape._element.find(qn("p:style"))
    if style is not None:
        shape._element.remove(style)


def _style_title(placeholder, size=None, color=ACCENT, anchor=MSO_ANCHOR.BOTTOM) -> None:
    """Report titles read left-aligned. The 4:3 master centres them, which on a 16:9
    canvas parks the title off-centre over the dead right-hand strip. The Section
    Header master also upper-cases its title, which mangles Vietnamese diacritics and
    matches no other slide in the deck."""
    frame = placeholder.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    chosen = size or _pick(TITLE_SIZES, len(placeholder.text))
    for para in frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT
        _paint(para, size=chosen, color=color, bold=True)
        para.font._rPr.set("cap", "none")
        for run in para.runs:
            run.font._rPr.set("cap", "none")


def _title_slide(prs: Presentation, title: str, subtitle: str,
                 cover_image: str | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    _color_background(slide, BRAND_BG)
    text_w = Inches(7.1) if cover_image else FULL_W
    slide.shapes.title.text = title
    _fit(slide.shapes.title, MARGIN, Inches(1.3), text_w, Inches(1.8))
    _style_title(slide.shapes.title, size=Pt(40), color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    sub = slide.placeholders[1]
    if subtitle:
        sub.text = subtitle
        _fit(sub, MARGIN, Inches(3.4), text_w, Inches(1.0))
        for para in sub.text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            _paint(para, size=Pt(18), color=TINT)
    else:
        # An empty placeholder still prints its "Click to add subtitle" prompt in
        # some viewers; drop the shape when there is nothing to say.
        sub._element.getparent().remove(sub._element)
    _accent_bar(slide, top=Inches(3.25), width=text_w, color=WHITE)
    if cover_image:
        _place_picture_cover(slide, cover_image, Inches(8.15), Inches(1.1), Inches(4.4), Inches(5.3))


def _is_divider(sec: dict) -> bool:
    """One short line under a heading reads as a lead-in to a part, not as content."""
    if not sec["first"] or not sec["title"]:
        return False
    body = [item for item in sec["items"] if item.strip()]
    return len(body) <= 1 and sum(len(item) for item in body) <= DIVIDER_MAX_CHARS


def _divider_slide(prs: Presentation, title: str, lead: str, image: str | None,
                   color_key: str) -> bool:
    color = _color_for(color_key)
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
    image_placed = False
    if image:
        try:
            # Full-bleed photo tinted with the section colour, title on top. The slide's
            # own background is painted the same dark tint even though the photo hides
            # it: `_background_is_dark` reads the background fill, and without this the
            # darkest slides in the deck reported "light" and got a grey slide number
            # drawn on near-black.
            _color_background(slide, _darken(color, SCRIM_DARKEN))
            picture = _place_picture_cover(slide, image, 0, 0, SLIDE_W, SLIDE_H)
            scrim = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
            scrim.fill.solid()
            scrim.fill.fore_color.rgb = _darken(color, SCRIM_DARKEN)
            _set_fill_alpha(scrim, SCRIM_ALPHA)
            scrim.line.fill.background()
            scrim.shadow.inherit = False
            _strip_style_ref(scrim)
            _send_to_back(scrim)
            _send_to_back(picture)
            image_placed = True
        except Exception:
            _color_background(slide, color)
    else:
        _color_background(slide, color)

    text_w = FULL_W
    slide.shapes.title.text = title
    _fit(slide.shapes.title, MARGIN, TITLE_TOP, text_w, Inches(1.5))
    _style_title(slide.shapes.title, size=_pick(TITLE_SIZES, len(title)),
                 color=WHITE, anchor=MSO_ANCHOR.TOP)

    body = slide.placeholders[1]
    _fit(body, MARGIN, Inches(2.3), text_w, Inches(1.3))
    if lead:
        body.text = lead
        frame = body.text_frame
        frame.word_wrap = True
        for para in frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            # White, not TINT: at 16pt the tint measures 4.1:1 on the green divider and
            # 2.3:1 once the divider carries a photo. TINT is only safe on the deep navy
            # of the cover and closing panels.
            _paint(para, size=Pt(16), color=WHITE)
    else:
        body._element.getparent().remove(body._element)
    _accent_bar(slide, top=TITLE_TOP + Inches(1.65), width=text_w, color=WHITE)

    return image_placed


def _bullet_slide(prs: Presentation, title: str, bullets: list[str],
                  image: str | None = None, color: RGBColor = ACCENT) -> bool:
    """Returns whether the picture actually landed -- a format python-pptx refuses
    must not be counted as imagery the deck carries."""
    layout = LAYOUT_TWO_CONTENT if image else LAYOUT_TITLE_BODY
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    slide.shapes.title.text = title
    _fit(slide.shapes.title, MARGIN, TITLE_TOP, FULL_W, TITLE_H)
    _style_title(slide.shapes.title, color=color)
    # A rule under the title gives every content slide the same visual anchor.
    _accent_bar(slide, top=TITLE_TOP + TITLE_H + Inches(0.06), height=Pt(2), color=color)

    body = slide.placeholders[1]
    _fit(body, MARGIN, BODY_TOP, TEXT_W_WITH_IMAGE if image else FULL_W, BODY_H)
    frame = body.text_frame
    frame.word_wrap = True
    # A slide carrying a couple of short lines looks abandoned pinned to the top edge.
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE if len(bullets) <= 3 else MSO_ANCHOR.TOP
    # Let PowerPoint shrink the text if it still overflows the placeholder.
    _enable_shrink_on_overflow(frame)
    size = _fitting_body_size(bullets, body.width, body.height)
    for i, text in enumerate(bullets):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.level = 0
        para.font.size = size
        para.font.color.rgb = BODY
        para.space_after = Pt(8)
        span = _lead_span(text)
        if span:
            _add_run(para, text[:span[1]], size, BODY, bold=True)
            _add_run(para, text[span[1]:], size, BODY)
        else:
            _add_run(para, text, size, BODY)

    if not image:
        return False
    # The picture placeholder only reserves the space; a real picture shape replaces it
    # so the image can be scaled to its own aspect ratio instead of being stretched.
    picture_ph = slide.placeholders[2]
    picture_ph._element.getparent().remove(picture_ph._element)
    try:
        _place_picture_cover(slide, image, PIC_LEFT, BODY_TOP, PIC_W, BODY_H)
        return True
    except Exception:
        return False  # bad image must not break the deck; the caller reports it


def _place_picture_cover(slide, image: str, left, top, width, height):
    """Fill the box with the picture, centre-cropping so it covers without distorting."""
    picture = slide.shapes.add_picture(image, left, top, width, height)
    w0, h0 = picture.image.size
    if w0 / h0 > width / height:
        crop = 1 - (h0 * width) / (height * w0)
        picture.crop_left = crop / 2
        picture.crop_right = crop / 2
    else:
        crop = 1 - (w0 * height) / (width * h0)
        picture.crop_top = crop / 2
        picture.crop_bottom = crop / 2
    return picture


def _send_to_back(shape) -> None:
    sp_tree = shape._element.getparent()
    sp_tree.remove(shape._element)
    sp_tree.insert(2, shape._element)


def _set_fill_alpha(shape, opacity_pct: float) -> None:
    """Set the solid fill's opacity (0-100). Used to tint a photo behind a divider."""
    from pptx.oxml.ns import qn

    solid = shape._element.spPr.find(qn("a:solidFill"))
    srgb = solid.find(qn("a:srgbClr")) if solid is not None else None
    if srgb is None:
        return
    for alpha in srgb.findall(qn("a:alpha")):
        srgb.remove(alpha)
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity_pct * 1000))}))


def _paint(para, size=None, color: RGBColor | None = None, bold: bool | None = None) -> None:
    """Apply font properties to the paragraph AND to every run inside it.

    Setting only ``paragraph.font`` writes ``a:pPr/a:defRPr``. PowerPoint and
    LibreOffice resolve a bare run against that default, so white-on-navy titles looked
    right there -- but Google Slides' importer ignores the paragraph default and falls
    back to its own layout, which repainted every title in the theme's dark ink and
    dropped the font size with it. A run that carries its own colour survives all three.
    """
    if size is not None:
        para.font.size = size
    if color is not None:
        para.font.color.rgb = color
    if bold is not None:
        para.font.bold = bold
    for run in para.runs:
        if size is not None:
            run.font.size = size
        if color is not None:
            run.font.color.rgb = color
        if bold is not None:
            run.font.bold = bold


def _add_run(para, text: str, size, color: RGBColor, bold: bool = False) -> None:
    if not text:
        return
    run = para.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold


def _card_backdrop(slide, left, top, width, height, color: RGBColor):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = color
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    _clear_effect_ref(card)
    return card


def _stat_slide(prs: Presentation, title: str, items: list[str], color: RGBColor) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
    slide.shapes.title.text = title
    _fit(slide.shapes.title, MARGIN, TITLE_TOP, FULL_W, TITLE_H)
    _style_title(slide.shapes.title, color=color)
    _accent_bar(slide, top=TITLE_TOP + TITLE_H + Inches(0.06), height=Pt(2), color=color)
    slide.placeholders[1]._element.getparent().remove(slide.placeholders[1]._element)

    items = [i for i in items if i.strip()]
    gap = Inches(0.16)
    card_h = min((SLIDE_H - BODY_TOP - Inches(0.5) - gap * (len(items) - 1)) / len(items),
                 Inches(1.3))
    top = BODY_TOP
    for item in items:
        card = _card_backdrop(slide, MARGIN, top, FULL_W, card_h, color)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.08)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _enable_shrink_on_overflow(tf)
        span = _find_stat_span(item)
        head = tf.paragraphs[0]
        body_w = card.width - Inches(0.3) - Inches(0.15)
        if span:
            head.font.size = Pt(24)
            _add_run(head, item[span[0]:span[1]], Pt(24), color, bold=True)
            body = tf.add_paragraph()
            head_h = Emu(int(Pt(24) * LINE_HEIGHT_EM))
            avail = card_h - Inches(0.1) - Inches(0.08) - head_h
            body_size = _fitting_card_body_size(item, body_w, avail)
            body.font.size = body_size
            _add_run(body, item, body_size, BODY)
        else:
            head.font.size = Pt(15)
            _add_run(head, item, Pt(15), BODY)
        top += card_h + gap


def _card_slide(prs: Presentation, title: str, items: list[str], color: RGBColor) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
    slide.shapes.title.text = title
    _fit(slide.shapes.title, MARGIN, TITLE_TOP, FULL_W, TITLE_H)
    _style_title(slide.shapes.title, color=color)
    _accent_bar(slide, top=TITLE_TOP + TITLE_H + Inches(0.06), height=Pt(2), color=color)
    slide.placeholders[1]._element.getparent().remove(slide.placeholders[1]._element)

    items = [i for i in items if i.strip()]
    cols = 3 if len(items) in (3, 5, 6) else 2
    gap = Inches(0.24)
    card_w = (FULL_W - gap * (cols - 1)) / cols
    rows = math.ceil(len(items) / cols)
    card_h = min((SLIDE_H - BODY_TOP - Inches(0.5) - gap * (rows - 1)) / rows, Inches(1.3))
    for i, item in enumerate(items):
        r, c = divmod(i, cols)
        card = _card_backdrop(slide, MARGIN + c * (card_w + gap),
                              BODY_TOP + r * (card_h + gap), card_w, card_h, color)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.12)
        tf.margin_bottom = Inches(0.08)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        _enable_shrink_on_overflow(tf)
        num = tf.paragraphs[0]
        num.font.size = Pt(20)
        _add_run(num, str(i + 1), Pt(20), color, bold=True)
        text = tf.add_paragraph()
        body_w = card_w - Inches(0.2) - Inches(0.2)
        num_h = Emu(int(Pt(20) * LINE_HEIGHT_EM))
        avail = card_h - Inches(0.12) - Inches(0.08) - num_h
        body_size = _fitting_card_body_size(item, body_w, avail)
        text.font.size = body_size
        _add_run(text, item, body_size, BODY)


def _enable_shrink_on_overflow(frame) -> None:
    """python-pptx has no setter for the placeholder's shrink-text-on-overflow, and
    PowerPoint needs the element present to apply it."""
    from pptx.oxml.ns import qn

    body_pr = frame._txBody.find(qn("a:bodyPr"))
    if body_pr is None:
        return
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        existing = body_pr.find(qn(tag))
        if existing is not None:
            body_pr.remove(existing)
    body_pr.append(body_pr.makeelement(qn("a:normAutofit"), {}))


def _credits_slide(prs: Presentation, credits: list[dict], words: dict) -> None:
    """CC-BY images must name the creator. cc0/pdm need not, but crediting all is simpler."""
    lines = []
    for credit in credits:
        title = (credit.get("title") or credit.get("query") or "Untitled")[:60]
        creator = credit.get("creator") or "Unknown"
        licence = credit.get("license") or "CC"
        lines.append(f"{title} — {creator} ({licence}) · Openverse")

    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_BODY])
    slide.shapes.title.text = words["credits"]
    _fit(slide.shapes.title, MARGIN, TITLE_TOP, FULL_W, TITLE_H)
    _style_title(slide.shapes.title)
    _accent_bar(slide, top=TITLE_TOP + TITLE_H + Inches(0.06), height=Pt(2))
    body = slide.placeholders[1]
    _fit(body, MARGIN, BODY_TOP, FULL_W, BODY_H)
    frame = body.text_frame
    frame.word_wrap = True
    _enable_shrink_on_overflow(frame)
    for i, line in enumerate(lines):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.text = line
        _paint(para, size=Pt(14), color=MUTED)
        para.space_after = Pt(6)


def _closing_slide(prs: Presentation, title: str, words: dict) -> None:
    """Close on a navy panel that mirrors the cover, with the title at the top."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
    _color_background(slide, BRAND_BG)
    slide.shapes.title.text = words["closing"]
    _fit(slide.shapes.title, MARGIN, TITLE_TOP, FULL_W, Inches(1.5))
    _style_title(slide.shapes.title, size=Pt(40), color=WHITE, anchor=MSO_ANCHOR.TOP)
    body = slide.placeholders[1]
    body.text = title
    _fit(body, MARGIN, Inches(2.3), FULL_W, Inches(1.2))
    for para in body.text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT
        _paint(para, size=Pt(20), color=TINT)
    _accent_bar(slide, top=Inches(2.05), color=WHITE)


def _accent_bar(slide, top, left=None, width=None, height=Pt(4), color=ACCENT) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN if left is None else left,
        top,
        FULL_W if width is None else width,
        height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    # An empty <a:effectLst/> is not enough on its own: the shape style still points at
    # the theme's effect #2, and LibreOffice honours that reference, drawing a drop
    # shadow under what should be a flat rule. Point the reference at "no effect".
    bar.shadow.inherit = False
    _clear_effect_ref(bar)


def _add_slide_numbers(prs: Presentation) -> None:
    """A report deck is referred to by slide number; the title slide keeps none."""
    for index, slide in enumerate(prs.slides, 1):
        if index == 1:
            continue
        color = TINT if _background_is_dark(slide) else MUTED
        box = slide.shapes.add_textbox(SLIDE_W - Inches(1.1), SLIDE_H - Inches(0.6),
                                       Inches(0.7), Inches(0.35))
        para = box.text_frame.paragraphs[0]
        para.text = str(index)
        _paint(para, size=Pt(11), color=color)
