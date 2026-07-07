"""Input handling for doc-convert: download links, detect type, extract a structured outline.

The common intermediate representation is a dict:
    {"title": str, "blocks": [{"kind": "heading"|"para"|"bullet", "level": int, "text": str}]}
"""
from __future__ import annotations

import os
import re
import subprocess
import urllib.parse
import urllib.request

SUPPORTED_EXTS = {".docx", ".pptx", ".pdf", ".txt", ".md", ".markdown"}

_DRIVE_FILE_RE = re.compile(r"drive\.google\.com/(?:file/d/|open\?id=|uc\?.*id=)([\w-]+)")
_GDOC_RE = re.compile(r"docs\.google\.com/document/d/([\w-]+)")
_GSLIDES_RE = re.compile(r"docs\.google\.com/presentation/d/([\w-]+)")


class DocConvertError(Exception):
    """User-visible conversion error."""


def is_url(value: str) -> bool:
    return value.startswith("http://") or value.startswith("https://")


def resolve_download_url(url: str) -> tuple[str, str | None]:
    """Map Google links to direct-download/export URLs. Returns (url, forced_ext)."""
    m = _GDOC_RE.search(url)
    if m:
        return (f"https://docs.google.com/document/d/{m.group(1)}/export?format=docx", ".docx")
    m = _GSLIDES_RE.search(url)
    if m:
        return (f"https://docs.google.com/presentation/d/{m.group(1)}/export/pptx", ".pptx")
    m = _DRIVE_FILE_RE.search(url)
    if m:
        return (f"https://drive.google.com/uc?export=download&id={m.group(1)}", None)
    return (url, None)


def download(url: str, dest_dir: str, timeout: int = 120) -> str:
    """Download a (public) file with curl, following redirects. Returns local path."""
    os.makedirs(dest_dir, exist_ok=True)
    direct_url, forced_ext = resolve_download_url(url)
    name = os.path.basename(urllib.parse.urlparse(direct_url).path) or "input"
    name = re.sub(r"[^\w.-]", "_", name)
    if forced_ext and not name.endswith(forced_ext):
        name += forced_ext
    dest = os.path.join(dest_dir, name)
    cmd = ["curl", "-fsSL", "--max-time", str(timeout), "-o", dest, direct_url]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0 or not os.path.exists(dest) or os.path.getsize(dest) == 0:
        raise DocConvertError(
            f"Download failed ({proc.returncode}): {proc.stderr.strip()[:200] or 'empty file'}. "
            "If this is a private Google link, enable link sharing or upload the file directly."
        )
    # Google returns an HTML error page instead of the file for private links.
    with open(dest, "rb") as fh:
        head = fh.read(512).lstrip().lower()
    if head.startswith(b"<!doctype html") or head.startswith(b"<html"):
        raise DocConvertError(
            "The link returned a web page instead of a file - it is probably private. "
            "Enable 'Anyone with the link' sharing or upload the file directly."
        )
    return dest


def detect_ext(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext in SUPPORTED_EXTS:
        return ".markdown" == ext and ".md" or ext
    raise DocConvertError(f"Unsupported input type: {ext or '(no extension)'}. Supported: docx, pptx, pdf, txt, md.")


def extract(path: str) -> dict:
    ext = detect_ext(path)
    if ext == ".docx":
        doc = _extract_docx(path)
    elif ext == ".pptx":
        doc = _extract_pptx(path)
    elif ext == ".pdf":
        doc = _extract_pdf(path)
    else:
        doc = _extract_text(path, markdown=ext == ".md")
    doc["blocks"] = [b for b in doc["blocks"] if b["text"].strip()]
    if not doc["blocks"]:
        raise DocConvertError("No extractable text found (scanned/image-only files are not supported).")
    if not doc.get("title"):
        doc["title"] = os.path.splitext(os.path.basename(path))[0].replace("_", " ")
    return doc


def _extract_docx(path: str) -> dict:
    import docx  # python-docx

    d = docx.Document(path)
    title = ""
    blocks: list[dict] = []
    for p in d.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        style = (p.style.name or "").lower()
        if style.startswith("title") and not title:
            title = text
        elif style.startswith("heading"):
            level = int(re.sub(r"\D", "", style) or 1)
            blocks.append({"kind": "heading", "level": min(level, 3), "text": text})
        elif "list" in style:
            blocks.append({"kind": "bullet", "level": 1, "text": text})
        else:
            blocks.append({"kind": "para", "level": 0, "text": text})
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                blocks.append({"kind": "bullet", "level": 1, "text": " - ".join(cells)})
    return {"title": title, "blocks": blocks}


def _extract_pptx(path: str) -> dict:
    from pptx import Presentation

    prs = Presentation(path)
    title = ""
    blocks: list[dict] = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_title = ""
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            slide_title = slide.shapes.title.text.strip()
        if idx == 1 and slide_title and not title:
            title = slide_title
        blocks.append({"kind": "heading", "level": 1, "text": slide_title or f"Slide {idx}"})
        for shape in slide.shapes:
            if not shape.has_text_frame or shape is slide.shapes.title:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    blocks.append({"kind": "bullet", "level": max(1, para.level + 1), "text": text})
        notes = getattr(slide, "notes_slide", None)
        if slide.has_notes_slide and notes and notes.notes_text_frame.text.strip():
            blocks.append({"kind": "para", "level": 0, "text": notes.notes_text_frame.text.strip()})
    return {"title": title, "blocks": blocks}


def _extract_pdf(path: str) -> dict:
    from pypdf import PdfReader

    reader = PdfReader(path)
    blocks: list[dict] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        for chunk in re.split(r"\n\s*\n", text):
            chunk = " ".join(chunk.split())
            if chunk:
                blocks.append({"kind": "para", "level": 0, "text": chunk})
    title = blocks[0]["text"][:80] if blocks else ""
    return {"title": title, "blocks": blocks}


def _extract_text(path: str, markdown: bool) -> dict:
    with open(path, encoding="utf-8", errors="replace") as fh:
        lines = fh.read().splitlines()
    title = ""
    blocks: list[dict] = []
    para: list[str] = []

    def flush() -> None:
        if para:
            blocks.append({"kind": "para", "level": 0, "text": " ".join(para)})
            para.clear()

    for line in lines:
        stripped = line.strip()
        if not stripped:
            flush()
            continue
        m = re.match(r"^(#{1,6})\s+(.*)$", stripped) if markdown else None
        if m:
            flush()
            level = len(m.group(1))
            if level == 1 and not title:
                title = m.group(2).strip()
            else:
                blocks.append({"kind": "heading", "level": min(level, 3), "text": m.group(2).strip()})
        elif markdown and re.match(r"^[-*+]\s+", stripped):
            flush()
            blocks.append({"kind": "bullet", "level": 1, "text": re.sub(r"^[-*+]\s+", "", stripped)})
        else:
            para.append(stripped)
    flush()
    return {"title": title, "blocks": blocks}


def to_markdown(doc: dict) -> str:
    out = [f"# {doc['title']}", ""]
    for b in doc["blocks"]:
        if b["kind"] == "heading":
            out.extend(["#" * (b["level"] + 1) + " " + b["text"], ""])
        elif b["kind"] == "bullet":
            out.append("  " * (b["level"] - 1) + "- " + b["text"])
        else:
            out.extend([b["text"], ""])
    return "\n".join(out).strip() + "\n"


def outline_sections(doc: dict) -> list[dict]:
    """Group blocks into sections: [{"title": str, "items": [str]}]."""
    sections: list[dict] = []
    current = {"title": "", "items": []}
    for b in doc["blocks"]:
        if b["kind"] == "heading":
            if current["items"] or current["title"]:
                sections.append(current)
            current = {"title": b["text"], "items": []}
        elif b["kind"] == "bullet":
            current["items"].append(b["text"])
        else:  # split long paragraphs into slide-sized statements
            for sent in _sentences(b["text"]):
                current["items"].append(sent)
    if current["items"] or current["title"]:
        sections.append(current)
    return sections


def _sentences(text: str, max_len: int = 160) -> list[str]:
    parts = re.split(r"(?<=[.!?])\s+", text)
    out: list[str] = []
    buf = ""
    for p in parts:
        if len(buf) + len(p) + 1 <= max_len:
            buf = (buf + " " + p).strip()
        else:
            if buf:
                out.append(buf)
            buf = p if len(p) <= max_len else p[: max_len - 1] + "…"
    if buf:
        out.append(buf)
    return out
