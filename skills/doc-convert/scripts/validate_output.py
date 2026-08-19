#!/usr/bin/env python3
"""Check a converted file against renderer-independent invariants.

Rendering a .pptx through LibreOffice only proves what LibreOffice would draw. The
product is opened in PowerPoint, Word, Canva and Google Slides, so the checks here
avoid asking any one renderer's opinion:

  * geometry is read straight out of the OOXML (every renderer honours it);
  * text is measured with the real Calibri-metric font file, so the fit holds for
    PowerPoint's Calibri and LibreOffice's Carlito alike;
  * content is compared against the source document, which catches text that a
    converter dropped or mangled no matter how pretty the result looks;
  * a deck imported into Google Slides is read back through the Slides API, because
    Drive's importer re-maps placeholders and re-wraps text after we hand the file
    over -- a clean local .pptx proves nothing about the deck Google ends up serving.

Usage:
  python3 validate_output.py --file <out.pptx|out.docx|out.pdf> [--source <input>]
  python3 validate_output.py --google <slides-url-or-id> [--file <exported.pdf>]

Prints a JSON report. Exit code 1 when a check fails.
"""
from __future__ import annotations

import argparse
import json
import os
import re
import shutil
import subprocess
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import doc_io  # noqa: E402

# The deck pins Calibri. Carlito is its metric-compatible clone, so measuring with
# whichever is installed gives line breaks that match both renderers.
FONT_CANDIDATES = ("Calibri", "Carlito")
# Diacritics that separate real Vietnamese from mojibake or stripped accents.
VI_DIACRITICS = "ăâđêôơưàáảãạèéẻẽẹìíỉĩịòóỏõọùúủũụỳýỷỹỵ"
# Ignore text shorter than this when comparing coverage: slide furniture such as
# numbers and the "cont." suffix has no counterpart in the source.
MIN_COVERAGE_CHARS = 25
# Content reaching less than this fraction of the slide width means the layout is
# still sized for a narrower canvas than the deck actually uses.
CANVAS_USE_MIN = 0.82

EMU_PER_PT = 12700
# Google Slides pads a text box by 0.1" on each side.
GSLIDES_TEXT_INSET_EMU = 91440
# The importer rounds geometry; 0.1" of slack keeps that off the bounds report.
GSLIDES_BOUNDS_TOL_EMU = 91440
# Slides re-wraps with its own metric substitute for Calibri and keeps the deck's
# shrink-on-overflow setting, so a box measuring slightly over is normal and would
# only produce noise. Only a box that cannot plausibly hold its text is a fault.
GSLIDES_OVERFLOW_TOL = 1.20


# WCAG AA: body text needs 4.5:1 against its background, large text 3:1. A reader on a
# phone or a projector loses more contrast than a reader on this monitor, so the deck is
# held to the standard rather than to "looks fine here".
CONTRAST_MIN = 4.5
CONTRAST_MIN_LARGE = 3.0
LARGE_PT = 24.0
LARGE_BOLD_PT = 18.0
# A scrim tinting a photo is measured against white: the brightest picture the image
# search can return is the case that has to stay readable.
ASSUMED_PHOTO = (255, 255, 255)


def _luminance(rgb: tuple) -> float:
    def channel(value: float) -> float:
        value /= 255
        return value / 12.92 if value <= 0.03928 else ((value + 0.055) / 1.055) ** 2.4
    r, g, b = rgb
    return 0.2126 * channel(r) + 0.7152 * channel(g) + 0.0722 * channel(b)


def contrast_ratio(fg: tuple, bg: tuple) -> float:
    high, low = sorted((_luminance(fg), _luminance(bg)), reverse=True)
    return (high + 0.05) / (low + 0.05)


def _blend(fg: tuple, bg: tuple, alpha: float) -> tuple:
    return tuple(alpha * f + (1 - alpha) * b for f, b in zip(fg, bg))


def _as_rgb(color) -> tuple | None:
    """RGB triple of a python-pptx colour, or None for theme/inherited colours."""
    try:
        value = int(str(color.rgb), 16)
    except Exception:  # noqa: BLE001 - MSO_THEME_COLOR and None both land here
        return None
    return ((value >> 16) & 0xFF, (value >> 8) & 0xFF, value & 0xFF)


def _solid_fill(shape) -> tuple | None:
    """(rgb, alpha) of a shape's solid fill, or None when it has no literal fill."""
    from pptx.oxml.ns import qn

    try:
        if shape.fill.type != 1:  # MSO_FILL.SOLID
            return None
        rgb = _as_rgb(shape.fill.fore_color)
    except Exception:  # noqa: BLE001 - pictures and placeholders raise here
        return None
    if rgb is None:
        return None
    alpha = 1.0
    solid = shape._element.spPr.find(qn("a:solidFill")) if shape._element.spPr is not None else None
    srgb = solid.find(qn("a:srgbClr")) if solid is not None else None
    if srgb is not None:
        node = srgb.find(qn("a:alpha"))
        if node is not None:
            alpha = int(node.get("val")) / 100000
    return rgb, alpha


def _slide_background(slide, canvas_w, canvas_h) -> tuple:
    """The colour text actually sits on: a full-bleed scrim if the slide has one, else
    the slide's own fill, else white."""
    background = (255, 255, 255)
    try:
        if slide.background.fill.type == 1:
            background = _as_rgb(slide.background.fill.fore_color) or background
    except Exception:  # noqa: BLE001 - inherited backgrounds have no literal colour
        pass
    for shape in slide.shapes:
        if shape.left is None or shape.width is None:
            continue
        covers = (shape.left <= canvas_w * 0.02 and shape.top <= canvas_h * 0.02
                  and shape.width >= canvas_w * 0.9 and shape.height >= canvas_h * 0.9)
        if not covers:
            continue
        fill = _solid_fill(shape)
        if fill:
            # The scrim sits over a picture, so blend it against the brightest photo.
            background = _blend(fill[0], ASSUMED_PHOTO, fill[1])
    return background


def check_contrast(slide, canvas_w, canvas_h) -> list[dict]:
    """Every run of text must clear WCAG AA against whatever it sits on."""
    problems: list[dict] = []
    base = _slide_background(slide, canvas_w, canvas_h)
    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        own = _solid_fill(shape)
        background = _blend(own[0], base, own[1]) if own else base
        for para in shape.text_frame.paragraphs:
            runs = para.runs or []
            for run in runs:
                text = run.text.strip()
                if not text:
                    continue
                colour = _as_rgb(run.font.color)
                if colour is None:
                    # The run leans on an inherited colour. PowerPoint and LibreOffice
                    # resolve that against the paragraph default, but Google Slides'
                    # importer repaints from its own layout -- a white-on-navy title
                    # comes back black. Harmless on a light slide, fatal on a dark one.
                    if _luminance(background) < 0.5:
                        problems.append({
                            "check": "contrast",
                            "detail": (f"{text[:40]!r} has no explicit run colour on a dark "
                                       f"background; Google Slides will repaint it"),
                        })
                    continue
                size = run.font.size or para.font.size
                points = size.pt if size else 18.0
                bold = bool(run.font.bold or para.font.bold)
                large = points >= LARGE_PT or (bold and points >= LARGE_BOLD_PT)
                needed = CONTRAST_MIN_LARGE if large else CONTRAST_MIN
                ratio = contrast_ratio(colour, background)
                if ratio < needed:
                    problems.append({
                        "check": "contrast",
                        "detail": (f"{text[:40]!r} at {points:.0f}pt measures "
                                   f"{ratio:.2f}:1 on #{'%02X%02X%02X' % tuple(int(c) for c in background)}"
                                   f" (needs {needed})"),
                    })
    return problems


def _font_path() -> str | None:
    if not shutil.which("fc-match"):
        return None
    for family in FONT_CANDIDATES:
        try:
            path = subprocess.run(["fc-match", "-f", "%{file}", family],
                                  capture_output=True, text=True, timeout=15).stdout.strip()
        except Exception:  # noqa: BLE001 - measurement is best-effort
            return None
        if path and os.path.basename(path).lower().startswith(tuple(f.lower() for f in FONT_CANDIDATES)):
            return path
    return None


def _measure(text: str, pt: float, font_path: str) -> tuple[int, int]:
    """Return (width, height) of one line in points, using the real font."""
    from PIL import ImageFont

    # Render at 4x for sub-point accuracy, then scale back.
    font = ImageFont.truetype(font_path, int(pt * 4))
    left, top, right, bottom = font.getbbox(text)
    return (right - left) / 4, (bottom - top) / 4


def _wrapped_lines(text: str, pt: float, width_pt: float, font_path: str) -> int:
    """Greedy word wrap with real metrics -- what both renderers do."""
    words = text.split()
    if not words:
        return 1
    lines, current = 1, ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if _measure(candidate, pt, font_path)[0] <= width_pt or not current:
            current = candidate
        else:
            lines += 1
            current = word
    return lines


def check_pptx(path: str, font_path: str | None) -> list[dict]:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Inches

    problems: list[dict] = []
    prs = Presentation(path)
    canvas_w, canvas_h = prs.slide_width, prs.slide_height

    for index, slide in enumerate(prs.slides, 1):
        def fail(check: str, detail: str) -> None:
            problems.append({"slide": index, "check": check, "detail": detail})

        title = slide.shapes.title
        if title is None or not title.text_frame.text.strip():
            fail("title_placeholder", "slide has no title placeholder or it is empty")
        if slide.slide_layout.name == "Blank":
            fail("layout", "slide uses the Blank layout, so it carries no structure")

        for problem in check_contrast(slide, canvas_w, canvas_h):
            fail(problem["check"], problem["detail"])

        rules, texts, right_edges = [], [], []
        for shape in slide.shapes:
            if shape.left is None or shape.width is None:
                continue
            # Slide numbers and other furniture sit against the right margin and would
            # mask a body that never leaves the old 4:3 box.
            if shape.width > Inches(2):
                right_edges.append(shape.left + shape.width)
            if shape.left < 0 or shape.top < 0:
                fail("bounds", f"{shape.name} starts off-canvas")
            if shape.left + shape.width > canvas_w or shape.top + shape.height > canvas_h:
                fail("bounds", f"{shape.name} overruns the canvas")
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append((shape.name, shape.top, shape.top + shape.height))
            elif shape._element.find(qn("p:style")) is not None:
                rules.append((shape.name, shape.top, shape.top + shape.height))

        # A 4:3 layout dropped onto a 16:9 canvas keeps its old 10" geometry and leaves
        # the right-hand strip empty on every slide. Nothing overflows, so only the
        # unused width gives it away.
        if right_edges and max(right_edges) < canvas_w * CANVAS_USE_MIN:
            fail("canvas_use",
                 f"content stops at {max(right_edges) / canvas_w:.0%} of the slide width")

        for rule_name, top, bottom in rules:
            for text_name, text_top, text_bottom in texts:
                if min(bottom, text_bottom) - max(top, text_top) > 0:
                    fail("overlap", f"{rule_name} crosses {text_name}")

        # The expensive part: does the text actually fit its box in *any* renderer?
        if font_path:
            for shape in slide.shapes:
                if not shape.has_text_frame or not shape.text_frame.text.strip():
                    continue
                width_pt = Emu(shape.width).pt
                used = 0.0
                for para in shape.text_frame.paragraphs:
                    if not para.text.strip():
                        continue
                    pt = (para.font.size or Emu(0)).pt if para.font.size else 18.0
                    lines = _wrapped_lines(para.text, pt, width_pt * 0.94, font_path)
                    used += lines * pt * 1.22 + (para.space_after.pt if para.space_after else 0)
                if used > Emu(shape.height).pt * 1.02:
                    fail("overflow",
                         f"{shape.name}: text needs ~{used:.0f}pt in a "
                         f"{Emu(shape.height).pt:.0f}pt box")
    return problems


def check_google_slides(data: dict, font_path: str | None) -> list[dict]:
    """Check the deck Google actually holds after importing our .pptx.

    Drive's importer is Google's code: it re-maps placeholders onto Slides layouts,
    substitutes fonts and re-wraps every paragraph. A clean local .pptx therefore says
    nothing about what came out the other side, so the same fit and bounds invariants
    are re-run against the geometry Google reports back.
    """
    problems: list[dict] = []
    page_w = data.get("page", {}).get("width_emu") or 0
    page_h = data.get("page", {}).get("height_emu") or 0
    slides = data.get("slides") or []
    if not slides:
        return [{"check": "slides", "detail": "the presentation has no slides"}]

    for slide in slides:
        index = slide["index"] + 1

        def fail(check: str, detail: str, _i: int = index) -> None:
            problems.append({"slide": _i, "check": check, "detail": detail})

        if not slide["texts"] and not slide["images"]:
            fail("empty_slide", "slide carries neither text nor a picture")

        # An unset page background is Slides' default white.
        background = slide.get("background") or (255, 255, 255)
        dark_slide = _luminance(background) < 0.5
        for box in slide["texts"]:
            for run in box.get("runs") or []:
                label = f"{run['text'][:40]!r}"
                if run["rgb"] is None:
                    # Google repainted this run from its own layout. On a dark slide that
                    # means the theme's dark ink on a dark background -- the title that
                    # came out black over a navy photo divider.
                    if dark_slide:
                        fail("contrast", f"{label} lost its colour in the import and "
                                         f"inherits Google's theme on a dark slide")
                    continue
                points = run["font_pt"] or box.get("font_pt") or 18.0
                large = points >= LARGE_PT or (run["bold"] and points >= LARGE_BOLD_PT)
                needed = CONTRAST_MIN_LARGE if large else CONTRAST_MIN
                measured = contrast_ratio(run["rgb"], background)
                if measured < needed:
                    fail("contrast", f"{label} at {points:.0f}pt measures {measured:.2f}:1 "
                                     f"on #{'%02X%02X%02X' % background} (needs {needed})")

        for box in slide["texts"]:
            name = box.get("object_id") or "text box"
            right = box["left_emu"] + box["width_emu"]
            bottom = box["top_emu"] + box["height_emu"]
            if box["left_emu"] < -GSLIDES_BOUNDS_TOL_EMU or box["top_emu"] < -GSLIDES_BOUNDS_TOL_EMU:
                fail("bounds", f"{name} starts off-canvas")
            elif (page_w and right > page_w + GSLIDES_BOUNDS_TOL_EMU) or \
                 (page_h and bottom > page_h + GSLIDES_BOUNDS_TOL_EMU):
                fail("bounds", f"{name} overruns the canvas")

            if not font_path or box["width_emu"] <= 0 or box["height_emu"] <= 0:
                continue
            width_pt = max((box["width_emu"] - 2 * GSLIDES_TEXT_INSET_EMU) / EMU_PER_PT, 1.0)
            height_pt = box["height_emu"] / EMU_PER_PT
            used = 0.0
            runs = [r for r in (box.get("runs") or []) if (r.get("text") or "").strip()]
            if runs:
                # A hero-number card mixes a large figure run with a smaller body run in
                # one box; box["font_pt"] is the max, which over-counts the body to
                # several lines at the figure's size. Measure each run at its own size.
                for run in runs:
                    rpt = run.get("font_pt") or box.get("font_pt") or 14.0
                    used += _wrapped_lines(run["text"], rpt, width_pt, font_path) * rpt * 1.22
            else:
                # Slides writes a soft break as U+000B and a paragraph break as U+000A.
                pt = box["font_pt"] or 14.0
                for para in re.split(r"[\n\v]", box["text"]):
                    if not para.strip():
                        continue
                    used += _wrapped_lines(para, pt, width_pt, font_path) * pt * 1.22
            if used > height_pt * GSLIDES_OVERFLOW_TOL:
                fail("overflow", f"{name}: text needs ~{used:.0f}pt in a {height_pt:.0f}pt box")
    return problems


def check_docx(path: str) -> list[dict]:
    import docx

    problems: list[dict] = []
    d = docx.Document(path)
    used = {p.style.name for p in d.paragraphs if p.text.strip()}
    if "List Bullet" in used and "Normal" not in used:
        problems.append({"check": "prose", "detail":
                         "every paragraph is a bullet -- the document has no body text"})
    if not any(name.startswith("Heading") or name == "Title" for name in used):
        problems.append({"check": "headings", "detail": "no headings, so navigation is flat"})
    for para in d.paragraphs:
        for run in para.runs:
            if run.font.name and run.font.name not in ("Calibri", "Carlito"):
                problems.append({"check": "font", "detail":
                                 f"unexpected font {run.font.name!r} on {para.text[:40]!r}"})
                break
    return problems


def check_pdf(path: str) -> list[dict]:
    from pypdf import PdfReader

    problems: list[dict] = []
    reader = PdfReader(path)
    if not reader.pages:
        return [{"check": "pages", "detail": "PDF has no pages"}]

    fonts: set[str] = set()
    for page in reader.pages:
        for ref in (page.get("/Resources", {}).get("/Font", {}) or {}).values():
            base = str(ref.get_object().get("/BaseFont") or "")
            fonts.add(base)
            # A non-embedded font is named without the ABCDEF+ subset prefix.
            if "+" not in base:
                problems.append({"check": "font_embedding", "detail":
                                 f"{base} is not embedded; the viewer will substitute it"})
    if not fonts:
        problems.append({"check": "font_embedding", "detail": "no fonts found in the PDF"})

    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    if "�" in text:
        problems.append({"check": "encoding", "detail": "replacement characters in the text layer"})
    return problems


def check_coverage(out_text: str, source: str) -> list[dict]:
    """Every substantial line of the source must survive into the output."""
    doc = doc_io.extract(source)
    haystack = _normalise(out_text)
    missing = []
    for block in doc["blocks"]:
        text = block["text"].strip()
        if len(text) < MIN_COVERAGE_CHARS:
            continue
        # Slides split long paragraphs, so match on the opening clause.
        probe = _normalise(text)[:60]
        if probe and probe not in haystack:
            missing.append(text[:70])
    problems = [{"check": "coverage", "detail": f"missing from output: {t!r}"} for t in missing[:10]]
    if len(missing) > 10:
        problems.append({"check": "coverage", "detail": f"...and {len(missing) - 10} more"})
    return problems


def _normalise(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip().lower()


def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        from pptx import Presentation
        return "\n".join(
            para.text
            for slide in Presentation(path).slides
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs
        )
    if ext == ".docx":
        import docx
        return "\n".join(p.text for p in docx.Document(path).paragraphs)
    if ext == ".pdf":
        from pypdf import PdfReader
        return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
    return ""


def vietnamese_ratio(text: str) -> float:
    letters = [ch for ch in text.lower() if ch.isalpha()]
    if not letters:
        return 0.0
    return sum(1 for ch in letters if ch in VI_DIACRITICS) / len(letters)


def main() -> int:
    ap = argparse.ArgumentParser(description="Validate a converted document")
    ap.add_argument("--file", action="append", default=[],
                    help="Output file to check (repeatable)")
    ap.add_argument("--google", help="Google Slides URL or id to check after import")
    ap.add_argument("--source", help="Original input, enables the content-coverage check")
    args = ap.parse_args()
    if not args.file and not args.google:
        ap.error("pass --file and/or --google")

    font_path = _font_path()
    report = {"font_used_for_measurement": font_path, "files": [], "success": True}

    if args.google:
        entry: dict = {"target": args.google, "type": "gslides", "problems": []}
        try:
            import google_io

            data = google_io.inspect_presentation(args.google)
            entry["slides"] = len(data.get("slides") or [])
            entry["problems"] += check_google_slides(data, font_path)
        except Exception as err:  # noqa: BLE001 - keep the JSON contract for the agent
            entry["problems"].append({"check": "readback",
                                      "detail": f"{type(err).__name__}: {err}"})
        entry["ok"] = not entry["problems"]
        report["success"] &= entry["ok"]
        report["files"].append(entry)

    for path in args.file:
        ext = os.path.splitext(path)[1].lower()
        entry: dict = {"path": path, "type": ext, "problems": []}
        if not os.path.exists(path):
            entry["problems"].append({"check": "exists", "detail": "file not found"})
        elif ext == ".pptx":
            entry["problems"] += check_pptx(path, font_path)
        elif ext == ".docx":
            entry["problems"] += check_docx(path)
        elif ext == ".pdf":
            entry["problems"] += check_pdf(path)
        else:
            entry["problems"].append({"check": "type", "detail": f"unsupported: {ext}"})

        if os.path.exists(path) and ext in (".pptx", ".docx", ".pdf"):
            text = extract_text(path)
            entry["vietnamese_diacritic_ratio"] = round(vietnamese_ratio(text), 4)
            if args.source:
                entry["problems"] += check_coverage(text, args.source)

        entry["ok"] = not entry["problems"]
        report["success"] &= entry["ok"]
        report["files"].append(entry)

    if not font_path:
        report.setdefault("warnings", []).append(
            "Calibri/Carlito not found - text-fit checks were skipped")

    print(json.dumps(report, indent=2, ensure_ascii=False))
    return 0 if report["success"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
