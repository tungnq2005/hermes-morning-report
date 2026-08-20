"""Unit tests for the doc-convert skill. Run from the workspace root:
    python3 -m unittest discover skills/doc-convert/tests
"""
from __future__ import annotations

import base64
import contextlib
import io
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import unittest

SCRIPTS = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "scripts")
sys.path.insert(0, SCRIPTS)

import doc_io  # noqa: E402
import build_pptx  # noqa: E402
import convert  # noqa: E402
import google_io  # noqa: E402
import image_search  # noqa: E402
import validate_output  # noqa: E402

# A 1x1 PNG. Enough for python-pptx to embed without pulling in Pillow fixtures.
TINY_PNG = base64.b64decode(
    b"iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk"
    b"YPhfDwAChwGA60e6kgAAAABJRU5ErkJggg=="
)


# An empty directory that never gets credentials: pointing the skill at it keeps every
# subprocess test off Google, whatever the machine running the suite has authorized.
# One stable path rather than a fresh mkdtemp per run, so the suite leaves nothing behind.
NO_CREDS_DIR = os.path.join(tempfile.gettempdir(), "docconv-tests-no-google-creds")
os.makedirs(NO_CREDS_DIR, exist_ok=True)


def offline_env(*, search: bool = False) -> dict:
    """Environment for a hermetic convert.py run: no Drive, no image search.

    Google is the skill's renderer of record, so a test that forgot this would upload
    to the operator's real Drive on the VPS.
    """
    env = os.environ.copy()
    env["DOC_CONVERT_GCREDS_DIR"] = NO_CREDS_DIR
    if search:
        env.pop(image_search.DISABLE_ENV, None)
    else:
        env[image_search.DISABLE_ENV] = "1"
    return env


def write_tiny_png(path: str) -> str:
    with open(path, "wb") as fh:
        fh.write(TINY_PNG)
    return path


def make_sample_docx(path: str) -> None:
    import docx

    d = docx.Document()
    d.add_heading("Kế hoạch dự án AI", level=0)
    d.add_heading("Mục tiêu", level=1)
    d.add_paragraph("Tự động hoá bản tin buổi sáng cho khách hàng.")
    d.add_paragraph("Chuyển đổi tài liệu đa định dạng.", style="List Bullet")
    d.add_heading("Phạm vi", level=1)
    d.add_paragraph("Giai đoạn một tập trung vào Telegram.", style="List Bullet")
    d.add_paragraph("Giai đoạn hai tích hợp Google Workspace.", style="List Bullet")
    d.add_heading("Tiến độ", level=1)
    d.add_paragraph("Hoàn thành trong hai tuần kể từ khi ký hợp đồng.")
    d.save(path)


def make_sample_pptx(path: str) -> None:
    from pptx import Presentation

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Demo Deck"
    s1.placeholders[1].text = "Subtitle"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Points"
    s2.placeholders[1].text_frame.text = "First point"
    prs.save(path)


class ExtractTests(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="docconv-test-")

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_docx_extract_structure(self):
        path = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(path)
        doc = doc_io.extract(path)
        self.assertEqual(doc["title"], "Kế hoạch dự án AI")
        kinds = {b["kind"] for b in doc["blocks"]}
        self.assertIn("heading", kinds)
        self.assertIn("bullet", kinds)
        sections = doc_io.outline_sections(doc)
        self.assertGreaterEqual(len(sections), 3)

    def test_pptx_extract(self):
        path = os.path.join(self.tmp, "deck.pptx")
        make_sample_pptx(path)
        doc = doc_io.extract(path)
        self.assertEqual(doc["title"], "Demo Deck")
        texts = " ".join(b["text"] for b in doc["blocks"])
        self.assertIn("First point", texts)

    def test_markdown_roundtrip(self):
        md_path = os.path.join(self.tmp, "notes.md")
        with open(md_path, "w", encoding="utf-8") as fh:
            fh.write("# Tiêu đề\n\n## Phần một\n\n- ý một\n- ý hai\n\nĐoạn văn.\n")
        doc = doc_io.extract(md_path)
        out = doc_io.to_markdown(doc)
        self.assertIn("## Phần một", out)
        self.assertIn("- ý một", out)

    def test_markdown_inline_markers_never_reach_the_slides(self):
        """Slides hold plain strings, so `**` has nothing to become and used to be shown."""
        md_path = os.path.join(self.tmp, "styled.md")
        with open(md_path, "w", encoding="utf-8") as fh:
            fh.write(
                "# Hướng dẫn **nhanh**\n\n"
                "---\n\n"
                "## Cách dùng\n\n"
                "- **Đầu vào**: Word (`.docx`), *nghiêng*, ~~bỏ~~\n"
                "- Xem [tài liệu](https://example.org/doc) để biết thêm\n"
                "- ***rất quan trọng***\n\n"
                "> Đoạn trích dẫn có `mã` bên trong.\n"
            )
        doc = doc_io.extract(md_path)
        texts = [b["text"] for b in doc["blocks"]]
        joined = doc["title"] + " " + " ".join(texts)

        for marker in ("**", "~~", "`", "](", "---"):
            self.assertNotIn(marker, joined, f"raw {marker!r} survived extraction")
        self.assertEqual(doc["title"], "Hướng dẫn nhanh")
        self.assertIn("Đầu vào: Word (.docx), nghiêng, bỏ", texts)
        self.assertIn("Xem tài liệu để biết thêm", texts)
        self.assertIn("rất quan trọng", texts)
        self.assertIn("Đoạn trích dẫn có mã bên trong.", texts)
        # the rule is a separator, not a bullet and not a paragraph of its own
        self.assertNotIn("---", texts)

    def test_unsupported_ext(self):
        path = os.path.join(self.tmp, "x.xyz")
        with open(path, "w") as fh:
            fh.write("data")
        with self.assertRaises(doc_io.DocConvertError):
            doc_io.extract(path)

    def test_google_link_resolution(self):
        url, ext = doc_io.resolve_download_url("https://docs.google.com/document/d/abc123XYZ/edit")
        self.assertIn("export?format=docx", url)
        self.assertEqual(ext, ".docx")
        url, ext = doc_io.resolve_download_url("https://docs.google.com/presentation/d/p456/edit#slide=1")
        self.assertIn("/export/pptx", url)


class BuildTests(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="docconv-test-")

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_pptx_min_slides(self):
        path = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(path)
        doc = doc_io.extract(path)
        sections = doc_io.outline_sections(doc)
        out = os.path.join(self.tmp, "out.pptx")
        stats = build_pptx.build(doc, sections, out, min_slides=5)
        self.assertTrue(os.path.exists(out))
        self.assertGreaterEqual(stats["slides"], 5)

    def test_built_docx_pins_font_and_bullet_glyph(self):
        """Guard the two defects that broke Vietnamese in the docx and pdf targets.

        Body text used to inherit the theme's minor font (Cambria), which LibreOffice
        replaced with a serif face lacking precomposed Vietnamese, and bullets used the
        Symbol-only codepoint U+F0B7, which rendered as a tofu box.
        """
        import zipfile

        import convert

        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        doc = doc_io.extract(src)
        out = os.path.join(self.tmp, "built.docx")
        convert.build_docx(doc, out)

        with zipfile.ZipFile(out) as z:
            styles = z.read("word/styles.xml").decode("utf-8")
            numbering = z.read("word/numbering.xml").decode("utf-8")

        for style_id in ("Normal", "ListBullet", "Heading1", "Heading2", "Title"):
            match = re.search(r'<w:style [^>]*w:styleId="%s".*?</w:style>' % style_id, styles, re.S)
            self.assertIsNotNone(match, f"style {style_id} missing")
            rfonts = re.search(r"<w:rFonts([^/>]*)/>", match.group(0))
            self.assertIsNotNone(rfonts, f"style {style_id} has no rFonts")
            self.assertIn('w:ascii="Calibri"', rfonts.group(1), f"style {style_id} not pinned")
            # a surviving theme reference would win over the literal font
            self.assertNotIn("Theme=", rfonts.group(1), f"style {style_id} still points at the theme")

        self.assertNotIn(convert.SYMBOL_BULLET, numbering, "Symbol-font bullet survived")
        self.assertIn(convert.UNICODE_BULLET, numbering, "U+2022 bullet missing")
        self.assertNotIn('w:ascii="Symbol"', numbering, "numbering still asks for the Symbol font")

    def test_built_docx_keeps_prose_out_of_the_bullet_list(self):
        """build_docx used to reuse the slide outline, which turns every paragraph into
        sentence-sized bullets. The Word file then contained no body text at all."""
        import convert

        md = os.path.join(self.tmp, "mixed.md")
        with open(md, "w", encoding="utf-8") as fh:
            fh.write(
                "# Báo cáo\n\n"
                "## Phần một\n\n"
                "Đây là một đoạn văn xuôi bình thường. Nó có hai câu.\n\n"
                "- ý gạch đầu dòng một\n"
                "- ý gạch đầu dòng hai\n"
            )
        out = os.path.join(self.tmp, "mixed.docx")
        convert.build_docx(doc_io.extract(md), out)

        import docx
        styled = [(p.style.name, p.text) for p in docx.Document(out).paragraphs if p.text.strip()]
        by_style = {}
        for style, text in styled:
            by_style.setdefault(style, []).append(text)

        self.assertIn("Normal", by_style, "prose must stay body text, not become a bullet")
        self.assertEqual(by_style["Normal"],
                         ["Đây là một đoạn văn xuôi bình thường. Nó có hai câu."],
                         "the paragraph must survive whole, not split per sentence")
        self.assertEqual(by_style["List Bullet"],
                         ["ý gạch đầu dòng một", "ý gạch đầu dòng hai"])
        # `##` is a level-2 heading and must stay one; the old code flattened every
        # section title to Heading 1.
        self.assertEqual(by_style["Heading 2"], ["Phần một"])

    def test_convert_cli_docx_to_pptx(self):
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        outdir = os.path.join(self.tmp, "run")
        proc = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "convert.py"),
             "--input", src, "--to", "pptx", "--outdir", outdir],
            capture_output=True, text=True, env=offline_env())
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        manifest = json.loads(proc.stdout)
        self.assertTrue(manifest["success"])
        self.assertTrue(os.path.exists(manifest["output"]))
        self.assertGreaterEqual(manifest["slides"], 5)

    def test_vietnamese_deck_refuses_to_guess_an_image(self):
        """Openverse answers Vietnamese queries with fishing boats. Better to ship no picture."""
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)  # Vietnamese
        outdir = os.path.join(self.tmp, "run-vi")
        # search is *allowed*; it must decline anyway
        proc = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "convert.py"),
             "--input", src, "--to", "pptx", "--outdir", outdir],
            capture_output=True, text=True, env=offline_env(search=True))
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        manifest = json.loads(proc.stdout)
        self.assertEqual(manifest["images_used"], 0)
        self.assertIn("image_search_needs_english_query", manifest["warnings"])
        self.assertNotIn("image_credits", manifest)

    def test_no_auto_images_flag_skips_search(self):
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        outdir = os.path.join(self.tmp, "run-noimg")
        proc = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "convert.py"),
             "--input", src, "--to", "pptx", "--outdir", outdir,
             "--image-query", "artificial intelligence", "--no-auto-images"],
            capture_output=True, text=True, env=offline_env(search=True))
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        manifest = json.loads(proc.stdout)
        self.assertEqual(manifest["images_used"], 0)
        # The only warning left is the one about rendering locally instead of in Google.
        self.assertEqual(manifest["warnings"], ["google_unauthorized:rendered_locally"])

    # ── Pictures in documents (gdoc / docx) ────────────────────────────
    def _docx_picture_layout(self, path: str) -> list[str]:
        """Read a .docx back as a flat list: heading/paragraph texts and 'PICTURE'."""
        import docx

        out = []
        for p in docx.Document(path).paragraphs:
            if "graphicData" in p._p.xml:
                out.append("PICTURE")
            elif p.text.strip():
                out.append(p.text.strip())
        return out

    def test_docx_puts_one_picture_under_each_section_heading(self):
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        doc = doc_io.extract(src)
        sections = doc_io.outline_sections(doc)
        png = write_tiny_png(os.path.join(self.tmp, "pic.png"))
        out = os.path.join(self.tmp, "illustrated.docx")

        stats = convert.build_docx(doc, out, sections=sections,
                                   images=[png] * len(sections), credits=[])
        self.assertEqual(stats["images_used"], len(sections))
        layout = self._docx_picture_layout(out)
        for heading in ("Mục tiêu", "Phạm vi", "Tiến độ"):
            self.assertEqual(layout[layout.index(heading) + 1], "PICTURE",
                             f"no picture directly under {heading}: {layout}")

    def test_docx_leaves_the_right_section_bare_when_a_picture_is_missing(self):
        """A hole in the image list must skip that section, not shift the next one."""
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        doc = doc_io.extract(src)
        sections = doc_io.outline_sections(doc)
        png = write_tiny_png(os.path.join(self.tmp, "pic.png"))
        images = [png] * len(sections)
        images[sections.index(next(s for s in sections if s["title"] == "Phạm vi"))] = None
        out = os.path.join(self.tmp, "gap.docx")

        stats = convert.build_docx(doc, out, sections=sections, images=images, credits=[])
        self.assertEqual(stats["images_used"], len(sections) - 1)
        layout = self._docx_picture_layout(out)
        self.assertNotEqual(layout[layout.index("Phạm vi") + 1], "PICTURE", layout)
        self.assertEqual(layout[layout.index("Tiến độ") + 1], "PICTURE", layout)

    def test_docx_credits_its_pictures(self):
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        doc = doc_io.extract(src)
        sections = doc_io.outline_sections(doc)
        png = write_tiny_png(os.path.join(self.tmp, "pic.png"))
        out = os.path.join(self.tmp, "credited.docx")

        convert.build_docx(doc, out, sections=sections, images=[png] * len(sections),
                           credits=[{"title": "Gold bars", "creator": "Jane Doe", "license": "by"}])
        text = "\n".join(self._docx_picture_layout(out))
        self.assertIn("Nguồn ảnh", text)  # the sample document is Vietnamese
        self.assertIn("Jane Doe", text)

    def test_docx_without_pictures_carries_no_credits_section(self):
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        doc = doc_io.extract(src)
        out = os.path.join(self.tmp, "plain.docx")

        stats = convert.build_docx(doc, out, sections=doc_io.outline_sections(doc),
                                   images=[], credits=[{"creator": "Jane Doe"}])
        self.assertEqual(stats["images_used"], 0)
        self.assertNotIn("Jane Doe", "\n".join(self._docx_picture_layout(out)))

    def test_plain_conversion_to_docx_stays_free_of_stock_photos(self):
        """Nobody wants Openverse photography inside a contract they only asked to convert."""
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        outdir = os.path.join(self.tmp, "run-docx-plain")
        proc = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "convert.py"),
             "--input", src, "--to", "docx", "--rebuild", "--outdir", outdir],
            capture_output=True, text=True, env=offline_env(search=True))
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        manifest = json.loads(proc.stdout)
        self.assertEqual(manifest["images_used"], 0)
        self.assertNotIn("PICTURE", self._docx_picture_layout(manifest["output"]))

    def test_docx_takes_the_pictures_the_caller_supplies(self):
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        png = write_tiny_png(os.path.join(self.tmp, "pic.png"))
        outdir = os.path.join(self.tmp, "run-docx-img")
        proc = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "convert.py"),
             "--input", src, "--to", "docx", "--rebuild", "--outdir", outdir,
             "--image", png, "--image", png],
            capture_output=True, text=True, env=offline_env(search=False))
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        manifest = json.loads(proc.stdout)
        self.assertEqual(manifest["images_used"], 2)
        self.assertIn("PICTURE", self._docx_picture_layout(manifest["output"]))

    def test_pptx_keeps_image_slots_aligned_when_one_fetch_fails(self):
        """A missing picture must leave *its own* slide bare, not shift the next one."""
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        doc = doc_io.extract(src)
        sections = doc_io.outline_sections(doc)
        good = write_tiny_png(os.path.join(self.tmp, "ok.png"))
        out = os.path.join(self.tmp, "aligned.pptx")

        # section 0 -> no image, section 1 -> image, section 2 -> broken path
        stats = build_pptx.build(doc, sections, out, images=[None, good, "/nonexistent.png"])
        self.assertEqual(stats["images_used"], 1)

        from pptx import Presentation
        prs = Presentation(out)
        picture_slides = [
            i for i, slide in enumerate(prs.slides)
            if any(sh.shape_type == 13 for sh in slide.shapes)  # 13 == PICTURE
        ]
        # title(0) + agenda(1) + sections(2,3,4): only the second section carries a picture
        self.assertEqual(picture_slides, [3])

    def test_pptx_keeps_image_slots_aligned_when_a_section_spans_slides(self):
        """A section split across slides must not eat the next sections' pictures."""
        doc = {"title": "Guide", "blocks": [{"kind": "para", "text": "english sample text"}]}
        sections = [
            {"title": "Alpha", "items": [f"a{i}" for i in range(8)]},  # 8 items -> 2 slides
            {"title": "Beta", "items": ["b1"]},
            {"title": "Gamma", "items": ["c1"]},
        ]
        pngs = [write_tiny_png(os.path.join(self.tmp, f"{n}.png")) for n in "abc"]
        out = os.path.join(self.tmp, "spanning.pptx")
        stats = build_pptx.build(doc, sections, out, min_slides=1, images=pngs)

        from pptx import Presentation
        carried = {
            slide.shapes.title.text
            for slide in Presentation(out).slides
            if slide.shapes.title is not None
            and any(sh.shape_type == 13 for sh in slide.shapes)  # 13 == PICTURE
        }
        self.assertEqual(carried, {"Alpha", "Beta", "Gamma"})
        self.assertNotIn("Alpha (cont.)", carried, "continuation slide stole the next picture")
        self.assertEqual(stats["images_used"], 3)

    def test_cover_image_is_not_repeated_inside_the_deck(self):
        doc = {"title": "Guide", "blocks": [{"kind": "para", "text": "english sample text"}]}
        sections = [{"title": "Alpha", "items": ["a1"]}, {"title": "Beta", "items": ["b1", "b2", "b3"]}]
        pngs = [write_tiny_png(os.path.join(self.tmp, f"{n}.png")) for n in "ab"]
        out = os.path.join(self.tmp, "cover.pptx")

        stats = build_pptx.build(doc, sections, out, min_slides=1, images=pngs, cover_image=pngs[0])
        from pptx import Presentation
        prs = Presentation(out)
        with_picture = [i for i, slide in enumerate(prs.slides)
                        if any(sh.shape_type == 13 for sh in slide.shapes)]
        self.assertIn(0, with_picture, "the cover slide has no picture")
        self.assertEqual(stats["images_used"], 2, "cover + Beta, with Alpha's slot released")

    def test_stat_heavy_deck_still_shows_a_picture(self):
        """Stat cards leave no room for a photo, so every fetched image used to vanish."""
        src = os.path.join(self.tmp, "report.md")
        with open(src, "w", encoding="utf-8") as fh:
            fh.write("# Gold market brief\n\n## Key moves\n"
                     "- Spot gold rose **1.2%** to **2,510 USD/ounce**.\n"
                     "- ETF inflows reached **480 million USD**.\n"
                     "- The dollar index fell **0.4%**.\n\n"
                     "## Why\n- US inflation came in at **2.6%** versus **2.8%** expected.\n")
        png = write_tiny_png(os.path.join(self.tmp, "gold.png"))
        outdir = os.path.join(self.tmp, "run-stat")
        proc = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "convert.py"),
             "--input", src, "--to", "pptx", "--outdir", outdir,
             "--image", png, "--image", png],
            capture_output=True, text=True, env=offline_env(search=False))
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        manifest = json.loads(proc.stdout)
        self.assertGreaterEqual(manifest["images_used"], 1,
                                "a number-heavy deck shipped with no imagery at all")

    def test_docx_recovers_a_jpeg_python_docx_cannot_parse(self):
        """Flickr serves JPEGs whose first marker is APP13; python-docx rejects those."""
        from PIL import Image

        raw = os.path.join(self.tmp, "photoshop.jpg")
        Image.new("RGB", (40, 30), (200, 160, 40)).save(raw, "JPEG")
        with open(raw, "rb") as fh:
            data = bytearray(fh.read())
        data[2:4] = b"\xff\xed"  # APP0 (JFIF) -> APP13 (Photoshop resource block)
        with open(raw, "wb") as fh:
            fh.write(data)

        import docx
        from docx.shared import Inches
        with self.assertRaises(Exception):
            docx.Document().add_picture(raw, width=Inches(1))

        doc = {"title": "Brief", "blocks": [{"kind": "heading", "level": 1, "text": "Alpha"},
                                            {"kind": "para", "text": "body"}]}
        out = os.path.join(self.tmp, "recovered.docx")
        stats = convert.build_docx(doc, out, sections=[{"title": "Alpha", "items": ["body"]}],
                                   images=[raw], credits=[])
        self.assertEqual(stats["images_used"], 1, stats)
        self.assertEqual(stats["images_rejected"], [])

    def test_pptx_slides_carry_real_placeholders(self):
        """Free-floating textboxes left the outline view empty and broke theme changes."""
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        doc = doc_io.extract(src)
        out = os.path.join(self.tmp, "structured.pptx")
        build_pptx.build(doc, doc_io.outline_sections(doc), out)

        from pptx import Presentation
        prs = Presentation(out)
        for index, slide in enumerate(prs.slides, 1):
            self.assertIsNotNone(slide.shapes.title, f"slide {index} has no title placeholder")
            self.assertTrue(slide.shapes.title.text.strip(), f"slide {index} title is empty")
            self.assertNotEqual(slide.slide_layout.name, "Blank", f"slide {index} is unstructured")

    def test_pptx_content_spans_the_widescreen_canvas(self):
        """The template is 4:3. Widening slide_width alone leaves every placeholder in
        the old 10" box, stranding 3.8" of dead space down the right-hand edge."""
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches

        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        doc = doc_io.extract(src)
        out = os.path.join(self.tmp, "wide.pptx")
        png = write_tiny_png(os.path.join(self.tmp, "p.png"))
        sections = doc_io.outline_sections(doc)
        build_pptx.build(doc, sections, out, images=[png] * len(sections))

        prs = Presentation(out)
        slide_w = prs.slide_width
        for index, slide in enumerate(prs.slides, 1):
            edges = []
            for shape in slide.shapes:
                right = shape.left + shape.width
                self.assertLessEqual(right, slide_w,
                                     f"slide {index}: {shape.name} overruns the canvas")
                edges.append(right)
            # Content confined to the old 4:3 box stops at 9.5"; the canvas is 13.333",
            # so anything short of ~11" means the right-hand strip was left dead.
            self.assertGreater(max(edges), Inches(11),
                               f"slide {index}: right-hand strip is unused")
            title = slide.shapes.title
            self.assertIsNotNone(title)
            self.assertEqual(title.text_frame.paragraphs[0].alignment, PP_ALIGN.LEFT)

    def test_pptx_titles_keep_their_own_casing(self):
        """The Section Header master upper-cases its title, which mangles Vietnamese."""
        doc = {"title": "Báo cáo", "blocks": [{"kind": "para", "text": "nội dung tiếng Việt"}]}
        sections = [{"title": "Chuyển đổi tài liệu", "items": ["Một câu dẫn ngắn."]}]
        out = os.path.join(self.tmp, "caps.pptx")
        build_pptx.build(doc, sections, out, min_slides=1)

        from pptx import Presentation
        for slide in Presentation(out).slides:
            title = slide.shapes.title
            if title is None:
                continue
            for para in title.text_frame.paragraphs:
                self.assertEqual(para.font._rPr.get("cap"), "none",
                                 f"title {para.text!r} still inherits ALL CAPS")

    def test_pptx_rules_are_flat_not_shadowed(self):
        """An empty effectLst does not stop LibreOffice honouring the theme effectRef."""
        from pptx.oxml.ns import qn

        doc = {"title": "Report", "blocks": [{"kind": "para", "text": "english body text"}]}
        out = os.path.join(self.tmp, "flat.pptx")
        build_pptx.build(doc, [{"title": "Alpha", "items": ["a", "b", "c"]}], out, min_slides=1)

        from pptx import Presentation
        bars = 0
        for slide in Presentation(out).slides:
            for shape in slide.shapes:
                style = shape._element.find(qn("p:style"))
                if style is None:
                    continue
                effect_ref = style.find(qn("a:effectRef"))
                if effect_ref is not None:
                    bars += 1
                    self.assertEqual(effect_ref.get("idx"), "0",
                                     "accent rule still references a theme shadow")
        self.assertGreater(bars, 0, "expected at least one accent rule")

    def test_pptx_rules_never_cross_a_text_box(self):
        """The title box is bottom-anchored, so a rule placed at its lower edge strikes
        through the descenders instead of underlining the title."""
        from pptx import Presentation
        from pptx.oxml.ns import qn

        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        doc = doc_io.extract(src)
        out = os.path.join(self.tmp, "rules.pptx")
        build_pptx.build(doc, doc_io.outline_sections(doc), out, subtitle="phụ đề")

        for index, slide in enumerate(Presentation(out).slides, 1):
            rules, texts = [], []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append((shape.name, shape.top, shape.top + shape.height))
                elif shape._element.find(qn("p:style")) is not None:
                    rules.append((shape.top, shape.top + shape.height))
            for top, bottom in rules:
                for name, text_top, text_bottom in texts:
                    overlap = min(bottom, text_bottom) - max(top, text_top)
                    self.assertLessEqual(overlap, 0,
                                         f"slide {index}: rule crosses {name}")

    def test_pptx_pads_by_splitting_content_not_by_inventing_slides(self):
        """The deck used to append slides whose only bullet read "(bổ sung)"."""
        doc = {"title": "Kế hoạch", "blocks": [{"kind": "para", "text": "nội dung tiếng Việt ở đây"}]}
        sections = [{"title": "Phần một", "items": [f"ý {i}" for i in range(6)]}]
        out = os.path.join(self.tmp, "padded.pptx")
        stats = build_pptx.build(doc, sections, out, min_slides=7)

        from pptx import Presentation
        text = "\n".join(
            para.text
            for slide in Presentation(out).slides
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs
        )
        self.assertNotIn("(bổ sung)", text)
        self.assertGreaterEqual(stats["slides"], 7)

    def test_pptx_localises_its_own_wording(self):
        """An English deck used to close on the Vietnamese slide "Tóm tắt & Q&A"."""
        english = {"title": "Quarterly Report", "blocks": [{"kind": "para", "text": "english body"}]}
        out = os.path.join(self.tmp, "en.pptx")
        build_pptx.build(english, [{"title": "Alpha", "items": ["a", "b"]}], out, min_slides=1)

        from pptx import Presentation
        titles = [s.shapes.title.text for s in Presentation(out).slides if s.shapes.title]
        self.assertIn("Summary & Q&A", titles)
        self.assertIn("Agenda", titles)
        self.assertNotIn("Tóm tắt & Q&A", titles)

    def test_credits_slide_names_creator_and_licence(self):
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        doc = doc_io.extract(src)
        sections = doc_io.outline_sections(doc)
        out = os.path.join(self.tmp, "credited.pptx")
        credits = [{"title": "Digital Economy", "creator": "ITU Pictures",
                    "license": "BY 2.0", "query": "digital economy"}]
        build_pptx.build(doc, sections, out, images=[write_tiny_png(os.path.join(self.tmp, "a.png"))],
                         credits=credits)

        from pptx import Presentation
        text = "\n".join(
            run.text
            for slide in Presentation(out).slides
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs for run in para.runs
        )
        self.assertIn("Nguồn ảnh", text)  # sample doc is Vietnamese
        self.assertIn("ITU Pictures", text)
        self.assertIn("BY 2.0", text)

    def test_convert_cli_pptx_to_docx(self):
        src = os.path.join(self.tmp, "deck.pptx")
        make_sample_pptx(src)
        outdir = os.path.join(self.tmp, "run2")
        proc = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "convert.py"),
             "--input", src, "--to", "docx", "--outdir", outdir],
            capture_output=True, text=True, env=offline_env())
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        manifest = json.loads(proc.stdout)
        self.assertTrue(manifest["success"])
        self.assertTrue(manifest["output"].endswith(".docx"))
        self.assertEqual(manifest["render_engine"], "local")

    @unittest.skipUnless(shutil.which("soffice"), "LibreOffice not installed")
    def test_convert_cli_docx_to_pdf(self):
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        outdir = os.path.join(self.tmp, "run3")
        proc = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "convert.py"),
             "--input", src, "--to", "pdf", "--outdir", outdir],
            capture_output=True, text=True, timeout=300, env=offline_env())
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        manifest = json.loads(proc.stdout)
        self.assertTrue(manifest["success"])
        self.assertTrue(manifest["output"].endswith(".pdf"))


class ImageSearchTests(unittest.TestCase):
    """No network. `_get` is the single seam through which image_search touches it."""

    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="imgsearch-test-")
        self.addCleanup(shutil.rmtree, self.tmp, True)
        self._real_get = image_search._get
        self.addCleanup(setattr, image_search, "_get", self._real_get)
        self.calls = []

    def fake_get(self, payload, content_type="application/json"):
        def _get(url, timeout):
            self.calls.append(url)
            return payload, content_type
        image_search._get = _get

    def test_disabled_by_env(self):
        os.environ[image_search.DISABLE_ENV] = "1"
        self.addCleanup(os.environ.pop, image_search.DISABLE_ENV, None)
        paths, credits, warnings = image_search.fetch(["digital economy"], self.tmp)
        self.assertEqual((paths, credits), ([], []))
        self.assertEqual(warnings, ["image_search_disabled"])
        self.assertEqual(self.calls, [], "disabled search must not touch the network")

    def test_search_requests_only_reusable_licences(self):
        self.fake_get(json.dumps({"results": [{
            "url": "https://example.org/a.jpg", "title": "Digital Economy",
            "creator": "ITU Pictures", "license": "by", "license_version": "2.0",
            "foreign_landing_url": "https://example.org/page",
        }]}).encode())
        hits = image_search.search("digital economy")

        self.assertEqual(len(self.calls), 1)
        # by-sa is viral and nc/nd bar client use; only cc0/pdm/by may be requested.
        self.assertIn("license=cc0%2Cpdm%2Cby", self.calls[0])
        self.assertIn("q=digital+economy", self.calls[0])
        self.assertEqual(hits[0]["creator"], "ITU Pictures")
        self.assertEqual(hits[0]["license"], "BY 2.0")

    def test_fetch_leaves_a_hole_when_one_query_finds_nothing(self):
        empty = json.dumps({"results": []}).encode()
        hit = json.dumps({"results": [{"url": "https://example.org/a.png", "title": "T",
                                       "creator": "C", "license": "cc0", "license_version": "1.0"}]}).encode()
        responses = [empty, hit]

        def _get(url, timeout):
            if url.startswith(image_search.API_URL):
                return responses.pop(0), "application/json"
            return TINY_PNG, "image/png"
        image_search._get = _get

        paths, credits, warnings = image_search.fetch(["nothing here", "cloud data center"], self.tmp)
        self.assertEqual(len(paths), 2, "one slot per query keeps slides aligned")
        self.assertIsNone(paths[0])
        self.assertTrue(paths[1].endswith(".png"))
        self.assertEqual(len(credits), 1)
        self.assertEqual(warnings, ["image_search_no_result:nothing here"])

    def test_search_prefers_titles_that_match_the_query(self):
        """Openverse's own order put a NASA belly-camera shot first for 'cloud data center'."""
        self.fake_get(json.dumps({"results": [
            {"url": "https://e.org/nasa.jpg", "title": "Global Hawk Pacific Belly Camera",
             "creator": "NASA", "license": "by", "license_version": "2.0"},
            {"url": "https://e.org/dc.jpg", "title": "Inside a cloud data center",
             "creator": "Someone", "license": "cc0", "license_version": "1.0"},
        ]}).encode())
        hits = image_search.search("cloud data center", limit=1)
        self.assertEqual(hits[0]["url"], "https://e.org/dc.jpg")

    def test_relevance_ignores_short_words(self):
        self.assertEqual(image_search.relevance("Inside a cloud data center", "cloud data center"), 3)
        self.assertEqual(image_search.relevance("A belly camera", "cloud data center"), 0)
        self.assertEqual(image_search.relevance("anything", "of in a"), 0)

    def test_download_retries_a_slow_origin_once(self):
        attempts = []

        def _get(url, timeout):
            attempts.append(url)
            if len(attempts) == 1:
                raise TimeoutError("read timed out")
            return TINY_PNG, "image/png"
        image_search._get = _get

        path = image_search.download({"url": "https://e.org/a.png"}, self.tmp, 1)
        self.assertEqual(len(attempts), 2)
        self.assertTrue(os.path.exists(path))

    def test_download_gives_up_after_the_retry(self):
        def _get(url, timeout):
            raise TimeoutError("read timed out")
        image_search._get = _get
        with self.assertRaises(image_search.ImageSearchError):
            image_search.download({"url": "https://e.org/a.png"}, self.tmp, 1)

    def test_download_rejects_webp_served_as_jpeg(self):
        """rawpixel serves WebP from a `.jpg` URL; python-pptx cannot embed it."""
        webp = b"RIFF\x24\x00\x00\x00WEBPVP8 "
        self.fake_get(webp, "image/webp")
        with self.assertRaises(image_search.UnsupportedImageFormat):
            image_search.download({"url": "https://example.org/a.jpg"}, self.tmp, 1)
        self.assertEqual(os.listdir(self.tmp), [], "unusable bytes must not be saved")

    def test_fetch_falls_through_to_the_next_candidate(self):
        """An unusable top hit costs the slide its picture only if every hit fails."""
        listing = json.dumps({"results": [
            {"url": "https://e.org/a.jpg", "title": "webp one", "creator": "A",
             "license": "cc0", "license_version": "1.0"},
            {"url": "https://e.org/b.png", "title": "good one", "creator": "B",
             "license": "cc0", "license_version": "1.0"},
        ]}).encode()

        def _get(url, timeout):
            if url.startswith(image_search.API_URL):
                return listing, "application/json"
            if url.endswith("a.jpg"):
                return b"RIFF\x24\x00\x00\x00WEBPVP8 ", "image/webp"
            return TINY_PNG, "image/png"
        image_search._get = _get

        paths, credits, warnings = image_search.fetch(["anything"], self.tmp)
        self.assertTrue(paths[0] and paths[0].endswith(".png"))
        self.assertEqual([c["creator"] for c in credits], ["B"], "credit the image we kept")
        self.assertEqual(warnings, ["image_unsupported_format:anything"])

    def test_download_rejects_a_non_image_body(self):
        self.fake_get(b"<html>404</html>", "text/html")
        with self.assertRaises(image_search.ImageSearchError):
            image_search.download({"url": "https://example.org/x"}, self.tmp, 1)

    def test_fetch_survives_a_dead_api(self):
        def _get(url, timeout):
            raise OSError("connection refused")
        image_search._get = _get
        paths, credits, warnings = image_search.fetch(["digital economy"], self.tmp)
        self.assertEqual(paths, [None])
        self.assertEqual(credits, [])
        self.assertEqual(warnings, ["image_search_failed:digital economy:OSError"])

    @unittest.skipUnless(os.environ.get("DOC_CONVERT_LIVE_TESTS") == "1",
                         "set DOC_CONVERT_LIVE_TESTS=1 to hit the real Openverse API")
    def test_live_openverse_returns_a_usable_image(self):
        paths, credits, warnings = image_search.fetch(["digital economy"], self.tmp)
        # A flaky origin may still warn; what must hold is that an image landed.
        self.assertTrue(paths and paths[0], f"no image came back: {warnings}")
        self.assertGreater(os.path.getsize(paths[0]), 1000)
        self.assertTrue(credits[0]["license"], "a credited licence is required by CC-BY")


class ValidateOutputTests(unittest.TestCase):
    """The validator is the only guard that survives a renderer disagreeing, so it has
    to be shown biting on the defects it exists to catch -- not just passing."""

    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="validate-test-")
        self.addCleanup(shutil.rmtree, self.tmp, True)

    def run_validator(self, *paths, source=None):
        cmd = [sys.executable, os.path.join(SCRIPTS, "validate_output.py")]
        for path in paths:
            cmd += ["--file", path]
        if source:
            cmd += ["--source", source]
        proc = subprocess.run(cmd, capture_output=True, text=True)
        return json.loads(proc.stdout)

    def checks_for(self, report, index=0):
        return {p["check"] for p in report["files"][index]["problems"]}

    def test_accepts_a_deck_the_builder_produced(self):
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        doc = doc_io.extract(src)
        out = os.path.join(self.tmp, "good.pptx")
        build_pptx.build(doc, doc_io.outline_sections(doc), out, subtitle="phụ đề")

        report = self.run_validator(out)
        self.assertTrue(report["success"], report["files"][0]["problems"])

    def test_flags_hand_placed_textboxes_and_a_narrow_canvas(self):
        """Reproduces the original builder: Blank layouts, no titles, 4:3 geometry."""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        for _ in range(2):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9.0), Inches(1.0))
            box.text_frame.paragraphs[0].text = "Tiêu đề vẽ tay"
            box.text_frame.paragraphs[0].font.size = Pt(28)
        out = os.path.join(self.tmp, "legacy.pptx")
        prs.save(out)

        checks = self.checks_for(self.run_validator(out))
        self.assertIn("title_placeholder", checks)
        self.assertIn("layout", checks)
        self.assertIn("canvas_use", checks)

    def test_flags_text_that_cannot_fit_its_box(self):
        """Measured with the real font, so the verdict holds outside LibreOffice."""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Tiêu đề"
        slide.shapes.title.left, slide.shapes.title.width = Inches(0.7), Inches(12)
        body = slide.placeholders[1]
        body.left, body.top, body.width, body.height = (
            Inches(0.7), Inches(1.6), Inches(11.9), Inches(1.0))
        frame = body.text_frame
        for i in range(8):
            para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            para.text = "Một dòng nội dung khá dài để chắc chắn tràn khỏi khung. " * 2
            para.font.size = Pt(20)
        out = os.path.join(self.tmp, "overflow.pptx")
        prs.save(out)

        self.assertIn("overflow", self.checks_for(self.run_validator(out)))

    def test_contrast_ratio_matches_the_wcag_reference_values(self):
        self.assertAlmostEqual(validate_output.contrast_ratio((0, 0, 0), (255, 255, 255)), 21.0, 2)
        self.assertAlmostEqual(validate_output.contrast_ratio((255, 255, 255), (255, 255, 255)), 1.0, 2)

    def test_every_run_stays_readable_even_over_the_brightest_photo(self):
        """A white picture under a divider's scrim is the worst case for white text;
        the deck must still clear WCAG AA there, not just on the flat slides."""
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        doc = doc_io.extract(src)
        sections = doc_io.outline_sections(doc)
        photo = write_tiny_png(os.path.join(self.tmp, "photo.png"))
        out = os.path.join(self.tmp, "contrast.pptx")
        build_pptx.build(doc, sections, out, images=[photo] * len(sections), subtitle="phụ đề")

        problems = [p for p in validate_output.check_pptx(out, None) if p["check"] == "contrast"]
        self.assertEqual(problems, [])

    def test_flags_text_that_disappears_into_its_background(self):
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        para = box.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = "chữ xám trên nền xám"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        out = os.path.join(self.tmp, "bad-contrast.pptx")
        prs.save(out)

        problems = [p for p in validate_output.check_pptx(out, None) if p["check"] == "contrast"]
        self.assertEqual(len(problems), 1, problems)
        self.assertIn("needs 4.5", problems[0]["detail"])

    def test_flags_a_run_that_leans_on_an_inherited_colour_over_a_dark_slide(self):
        """PowerPoint resolves a bare run against the paragraph default, so a white
        title looked right locally. Google Slides repaints it from its own layout, and
        the deck came back with black titles on the navy dividers."""
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0x0F, 0x2A, 0x4A)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        para = box.text_frame.paragraphs[0]
        para.font.size = Pt(28)
        para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # paragraph default only
        para.text = "tiêu đề trắng theo kiểu cũ"
        out = os.path.join(self.tmp, "inherited-colour.pptx")
        prs.save(out)

        problems = [p for p in validate_output.check_pptx(out, None) if p["check"] == "contrast"]
        self.assertEqual(len(problems), 1, problems)
        self.assertIn("no explicit run colour", problems[0]["detail"])

    def test_flags_a_word_file_with_no_body_text(self):
        import convert

        md = os.path.join(self.tmp, "prose.md")
        with open(md, "w", encoding="utf-8") as fh:
            fh.write("# Báo cáo\n\n## Phần\n\nĐoạn văn xuôi dài đủ để tính là nội dung thật.\n")
        good = os.path.join(self.tmp, "good.docx")
        convert.build_docx(doc_io.extract(md), good)
        self.assertTrue(self.run_validator(good)["success"])

        # Now the defect: every paragraph forced into the bullet style.
        import docx
        d = docx.Document()
        d.add_heading("Báo cáo", level=0)
        for text in ("Đoạn văn xuôi thứ nhất.", "Đoạn văn xuôi thứ hai."):
            d.add_paragraph(text, style="List Bullet")
        bad = os.path.join(self.tmp, "bad.docx")
        d.save(bad)
        self.assertIn("prose", self.checks_for(self.run_validator(bad)))

    def test_coverage_notices_content_dropped_on_the_floor(self):
        md = os.path.join(self.tmp, "source.md")
        with open(md, "w", encoding="utf-8") as fh:
            fh.write(
                "# Báo cáo\n\n"
                "Một đoạn văn đủ dài để được tính vào phần kiểm tra bao phủ nội dung.\n\n"
                "Một đoạn khác cũng đủ dài và bắt buộc phải xuất hiện trong bản xuất ra.\n"
            )
        import convert
        full = os.path.join(self.tmp, "full.docx")
        convert.build_docx(doc_io.extract(md), full)
        self.assertTrue(self.run_validator(full, source=md)["success"])

        truncated = doc_io.extract(md)
        truncated["blocks"] = truncated["blocks"][:1]
        partial = os.path.join(self.tmp, "partial.docx")
        convert.build_docx(truncated, partial)
        self.assertIn("coverage", self.checks_for(self.run_validator(partial, source=md)))


class NarrateTests(unittest.TestCase):
    def test_narrate_script(self):
        tmp = tempfile.mkdtemp(prefix="docconv-test-")
        self.addCleanup(shutil.rmtree, tmp, True)
        src = os.path.join(tmp, "sample.docx")
        make_sample_docx(src)
        proc = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "narrate.py"),
             "--input", src, "--lang", "Vietnamese", "--outdir", os.path.join(tmp, "aud")],
            capture_output=True, text=True)
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        result = json.loads(proc.stdout)
        self.assertTrue(result["success"])
        self.assertTrue(os.path.exists(result["script_path"]))
        self.assertIn("generate_audio.py", result["suggested_tts_command"])


class GoogleTests(unittest.TestCase):
    def test_url_detection_and_id(self):
        self.assertTrue(google_io.is_google_url("https://docs.google.com/document/d/ABC123def/edit"))
        self.assertTrue(google_io.is_google_url("https://docs.google.com/presentation/d/PID456/edit"))
        self.assertTrue(google_io.is_google_url("https://drive.google.com/file/d/DRV789/view"))
        self.assertFalse(google_io.is_google_url("https://example.com/x.docx"))
        self.assertEqual(google_io.extract_file_id(
            "https://docs.google.com/document/d/ABC123def/edit"), "ABC123def")
        self.assertEqual(google_io.extract_file_id(
            "https://drive.google.com/file/d/DRV789/view"), "DRV789")

    def test_scope_sets_stay_as_small_as_the_features_need(self):
        """Every restricted scope costs the customer a verification review, so the set
        must not carry scopes the code stopped calling."""
        self.assertEqual(google_io.SCOPE_SETS["minimal"], [google_io.SCOPE_DRIVE_FILE])
        self.assertNotIn("https://www.googleapis.com/auth/documents",
                         google_io.SCOPE_SETS["private-links"])
        self.assertNotIn("https://www.googleapis.com/auth/presentations",
                         google_io.SCOPE_SETS["private-links"])

    def test_scope_set_comes_from_the_environment_and_falls_back_safely(self):
        for value, expected in (("minimal", "minimal"), ("private-links", "private-links"),
                                ("nonsense", google_io.DEFAULT_SCOPE_SET)):
            os.environ["DOC_CONVERT_GOOGLE_SCOPES"] = value
            self.addCleanup(os.environ.pop, "DOC_CONVERT_GOOGLE_SCOPES", None)
            self.assertEqual(google_io.scope_set_name(), expected)

    def test_private_link_is_refused_with_an_actionable_message_on_a_minimal_token(self):
        tmp = tempfile.mkdtemp(prefix="gcreds-")
        self.addCleanup(shutil.rmtree, tmp, True)
        with open(os.path.join(tmp, "token.json"), "w", encoding="utf-8") as fh:
            json.dump({"scopes": [google_io.SCOPE_DRIVE_FILE]}, fh)

        self.assertFalse(google_io.can_read_private_files(tmp))
        with self.assertRaises(google_io.GoogleAuthError) as caught:
            google_io.download_private("https://docs.google.com/document/d/ABC123def/edit",
                                       tmp, creds_dir=tmp)
        self.assertIn("tải file lên", str(caught.exception))

    def test_has_token_false_in_empty_dir(self):
        tmp = tempfile.mkdtemp(prefix="gcreds-")
        self.addCleanup(shutil.rmtree, tmp, True)
        self.assertFalse(google_io.has_token(tmp))

    def test_load_credentials_without_token_raises(self):
        tmp = tempfile.mkdtemp(prefix="gcreds-")
        self.addCleanup(shutil.rmtree, tmp, True)
        with self.assertRaises(google_io.GoogleAuthError):
            google_io.load_credentials(tmp)

    def test_convert_cli_gdoc_without_token_fails_cleanly(self):
        # gdoc target must fail with an actionable message, not a crash, when unauthorized.
        tmp = tempfile.mkdtemp(prefix="docconv-test-")
        self.addCleanup(shutil.rmtree, tmp, True)
        src = os.path.join(tmp, "sample.docx")
        make_sample_docx(src)
        env = dict(os.environ)
        # Hermetic: point creds dir at an empty temp dir so we never touch real creds/API.
        env["DOC_CONVERT_GCREDS_DIR"] = os.path.join(tmp, "empty-creds")
        proc = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "convert.py"),
             "--input", src, "--to", "gdoc", "--outdir", os.path.join(tmp, "run")],
            capture_output=True, text=True, env=env)
        manifest = json.loads(proc.stdout)
        self.assertFalse(manifest["success"])
        self.assertIn("authorize", manifest["error"].lower())


def fake_deck(slides: int = 3) -> dict:
    """A geometrically clean deck as the Slides API would report it back."""
    return {
        "id": "PID",
        "page": {"width_emu": 12192000, "height_emu": 6858000},
        "slides": [
            {"index": i, "images": 0,
             "texts": [{"object_id": f"t{i}", "text": "Tiêu đề ngắn", "font_pt": 24.0,
                        "left_emu": 640080, "top_emu": 365760,
                        "width_emu": 10000000, "height_emu": 1000000}]}
            for i in range(slides)
        ],
    }


class GoogleRenderTests(unittest.TestCase):
    """Google renders the deliverable, so the pipeline is tested with the two Drive
    calls stubbed: no network, and no files created in anyone's real Drive."""

    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="grender-test-")
        self.addCleanup(shutil.rmtree, self.tmp, True)
        self.uploaded: list[tuple] = []
        self.exported: list[tuple] = []
        for name in ("has_token", "import_local", "export_to", "inspect_presentation"):
            self.addCleanup(setattr, google_io, name, getattr(google_io, name))
        google_io.has_token = lambda *a, **k: True
        google_io.import_local = self.fake_import
        google_io.export_to = self.fake_export
        google_io.inspect_presentation = lambda *a, **k: fake_deck()

    def fake_import(self, path, kind, title="", creds_dir=None):
        self.uploaded.append((path, kind, title))
        url = ("https://docs.google.com/presentation/d/PID/edit" if kind == "gslides"
               else "https://docs.google.com/document/d/DID/edit")
        return {"id": "PID" if kind == "gslides" else "DID", "url": url,
                "name": title, "kind": kind}

    def fake_export(self, file_id, fmt, dest_path, creds_dir=None):
        self.exported.append((file_id, fmt, dest_path))
        os.makedirs(os.path.dirname(dest_path), exist_ok=True)
        with open(dest_path, "wb") as fh:
            fh.write(b"GOOGLE-EXPORT")
        return dest_path

    def run_convert(self, *argv) -> dict:
        buf = io.StringIO()
        saved = sys.argv
        sys.argv = ["convert.py", *argv]
        try:
            with contextlib.redirect_stdout(buf):
                convert.main()
        finally:
            sys.argv = saved
        return json.loads(buf.getvalue())

    def test_gslides_uploads_the_built_deck_and_ships_a_pdf(self):
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        manifest = self.run_convert("--input", src, "--to", "gslides", "--no-auto-images",
                                    "--outdir", os.path.join(self.tmp, "run"))
        self.assertTrue(manifest["success"], manifest)
        path, kind, _title = self.uploaded[0]
        self.assertEqual(kind, "gslides")
        self.assertTrue(path.endswith(".pptx") and os.path.exists(path))
        self.assertEqual(manifest["google_url"], "https://docs.google.com/presentation/d/PID/edit")
        self.assertEqual(manifest["render_engine"], "google")
        self.assertTrue(manifest["output"].endswith(".pdf"))
        self.assertEqual(manifest["google_check"]["status"], "pass")

    def test_office_target_ships_googles_export_not_our_own_render(self):
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        manifest = self.run_convert("--input", src, "--to", "pptx", "--no-auto-images",
                                    "--outdir", os.path.join(self.tmp, "run"))
        self.assertEqual(manifest["render_engine"], "google")
        self.assertTrue(manifest["output"].endswith(".pptx"))
        with open(manifest["output"], "rb") as fh:
            self.assertEqual(fh.read(), b"GOOGLE-EXPORT")
        self.assertNotEqual(os.path.abspath(manifest["output"]),
                            os.path.abspath(manifest["local_build"]))

    def test_a_failed_export_still_delivers_the_link(self):
        def boom(*a, **k):
            raise google_io.GoogleExportError("Drive từ chối export pdf (giới hạn 10MB)")

        google_io.export_to = boom
        src = os.path.join(self.tmp, "sample.docx")
        make_sample_docx(src)
        manifest = self.run_convert("--input", src, "--to", "gslides", "--no-auto-images",
                                    "--outdir", os.path.join(self.tmp, "run"))
        self.assertTrue(manifest["success"])
        self.assertIn("google_url", manifest)
        self.assertNotIn("output", manifest)
        self.assertTrue(any(w.startswith("google_export_failed:") for w in manifest["warnings"]),
                        manifest["warnings"])

    def test_the_users_own_deck_is_uploaded_untouched_unless_rebuild_is_asked(self):
        src = os.path.join(self.tmp, "deck.pptx")
        make_sample_pptx(src)
        self.run_convert("--input", src, "--to", "gslides", "--no-auto-images",
                         "--outdir", os.path.join(self.tmp, "as-is"))
        self.assertIn(f"input{os.sep}deck.pptx", self.uploaded[0][0])

        self.run_convert("--input", src, "--to", "gslides", "--no-auto-images", "--rebuild",
                         "--outdir", os.path.join(self.tmp, "rebuilt"))
        self.assertIn(f"build{os.sep}deck.pptx", self.uploaded[1][0])


class GoogleDeckCheckTests(unittest.TestCase):
    """The readback is the only evidence about what Drive's importer produced."""

    def test_accepts_a_clean_deck(self):
        self.assertEqual(validate_output.check_google_slides(fake_deck(), None), [])

    def test_flags_an_empty_slide_and_a_box_off_the_canvas(self):
        data = fake_deck(2)
        data["slides"][0]["texts"] = []          # import dropped the text
        data["slides"][1]["texts"][0]["left_emu"] = 11000000  # box pushed off the page
        checks = {p["check"] for p in validate_output.check_google_slides(data, None)}
        self.assertEqual(checks, {"empty_slide", "bounds"})

    def test_flags_a_presentation_that_arrived_empty(self):
        data = fake_deck(0)
        problems = validate_output.check_google_slides(data, None)
        self.assertEqual([p["check"] for p in problems], ["slides"])


if __name__ == "__main__":
    unittest.main()
